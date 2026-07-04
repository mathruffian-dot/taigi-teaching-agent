# 自然語言教材產出資料夾驗證器
import json
import re
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List


REQUIRED_OUTPUT_KEYS = [
    "lesson_structure",
    "student_worksheet",
    "teacher_guide",
    "exam_paper",
    "exam_answer_key",
    "slides",
    "interactive_website",
    "teacher_review_report",
]


def read_json(path: Path) -> Any:
    with path.open("r", encoding="utf-8-sig") as f:
        return json.load(f)


def add_issue(issues: List[Dict[str, str]], severity: str, item: str, message: str) -> None:
    issues.append({"severity": severity, "item": item, "message": message})


def topic_terms(topic: str) -> List[str]:
    topic = (topic or "").strip()
    if not topic:
        return []
    terms: List[str] = []
    terms.extend(re.findall(r"[a-zA-Z0-9._-]+", topic.lower()))
    for chunk in re.findall(r"[\u4e00-\u9fff]+", topic):
        if len(chunk) <= 4:
            terms.append(chunk)
        if len(chunk) >= 4:
            terms.extend(chunk[idx:idx + 2] for idx in range(0, len(chunk) - 1))
    unique_terms = []
    for term in terms:
        if len(term) < 2:
            continue
        if term not in unique_terms:
            unique_terms.append(term)
    return unique_terms


def lesson_content_text(lesson: Dict[str, Any]) -> str:
    parts: List[str] = []
    for vocab in lesson.get("vocabulary", []) or []:
        if isinstance(vocab, dict):
            parts.extend(str(vocab.get(key, "")) for key in ("hanji", "tailo_numeric", "tailo_diacritic", "zh_tw"))
        else:
            parts.append(str(vocab))
    for dialogue in lesson.get("dialogues", []) or []:
        parts.extend(str(dialogue.get(key, "")) for key in ("role", "hanji", "tailo_numeric", "tailo_diacritic", "zh_tw"))
    for question in lesson.get("questions", []) or []:
        parts.extend(str(question.get(key, "")) for key in ("question", "explanation"))
        parts.extend(str(option) for option in question.get("options", []) or [])
    return "\n".join(parts)


def validate_topic_alignment(
    lesson: Dict[str, Any],
    manifest: Dict[str, Any],
    issues: List[Dict[str, str]]
) -> Dict[str, Any]:
    topic = (manifest.get("request") or {}).get("topic", "")
    terms = topic_terms(topic)
    content = lesson_content_text(lesson)
    matched_terms = [term for term in terms if term.lower() in content.lower()]
    if topic and terms and not matched_terms:
        add_issue(
            issues,
            "error",
            "topic_alignment",
            f"教材內容未出現自然語言主題關鍵詞：{topic}"
        )
    return {
        "topic": topic,
        "terms": terms,
        "matched_terms": matched_terms,
        "is_aligned": bool(matched_terms) if topic and terms else True,
    }


def validate_docx(path: Path, issues: List[Dict[str, str]], item: str) -> Dict[str, Any]:
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if not paragraphs:
            add_issue(issues, "error", item, "Word 文件沒有可讀文字。")
        return {"paragraph_count": len(paragraphs), "first_text": paragraphs[0] if paragraphs else ""}
    except Exception as exc:
        add_issue(issues, "error", item, f"Word 文件無法開啟：{exc}")
        return {}


def validate_pptx(path: Path, issues: List[Dict[str, str]]) -> Dict[str, Any]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        if len(prs.slides) < 5:
            add_issue(issues, "warning", "slides", f"簡報投影片數偏少：{len(prs.slides)}")
        return {"slide_count": len(prs.slides)}
    except Exception as exc:
        add_issue(issues, "error", "slides", f"PowerPoint 簡報無法開啟：{exc}")
        return {}


def validate_html(path: Path, issues: List[Dict[str, str]]) -> Dict[str, Any]:
    try:
        html = path.read_text(encoding="utf-8")
    except Exception as exc:
        add_issue(issues, "error", "interactive_website", f"HTML 無法讀取：{exc}")
        return {}

    required_markers = ["隨堂自我檢測", "核心詞彙", "情境會話"]
    missing = [marker for marker in required_markers if marker not in html]
    for marker in missing:
        add_issue(issues, "error", "interactive_website", f"HTML 缺少區塊：{marker}")
    return {
        "bytes": len(html.encode("utf-8")),
        "has_official_materials": "官方素材建議" in html,
    }


def validate_quiz_bank(path: Path, issues: List[Dict[str, str]]) -> Dict[str, Any]:
    try:
        quiz_bank = read_json(path)
    except Exception as exc:
        add_issue(issues, "error", "quiz_bank", f"測驗題庫 JSON 無法解析：{exc}")
        return {}

    questions = quiz_bank.get("questions", []) or []
    official_questions = quiz_bank.get("official_extension_questions", []) or []
    if not questions:
        add_issue(issues, "error", "quiz_bank", "測驗題庫缺少自動計分題。")
    for idx, question in enumerate(questions, start=1):
        options = question.get("options", []) or []
        answer_index = question.get("answer_index")
        if not question.get("question"):
            add_issue(issues, "error", "quiz_bank", f"第 {idx} 題缺少題幹。")
        if len(options) < 2:
            add_issue(issues, "error", "quiz_bank", f"第 {idx} 題選項少於 2 個。")
        if not isinstance(answer_index, int) or not 0 <= answer_index < len(options):
            add_issue(issues, "error", "quiz_bank", f"第 {idx} 題答案索引不合法。")
    return {
        "question_count": len(questions),
        "official_extension_count": len(official_questions),
    }


def resolve_output_path(output_dir: Path, value: str) -> Path:
    path = Path(value)
    if path.is_absolute():
        return path
    if path.exists():
        return path
    return output_dir / path.name


def validate_output_folder(output_dir: str) -> Dict[str, Any]:
    root = Path(output_dir)
    issues: List[Dict[str, str]] = []
    details: Dict[str, Any] = {}

    if not root.exists():
        add_issue(issues, "error", "output_dir", f"資料夾不存在：{root}")
        return {
            "validated_at": datetime.now().isoformat(timespec="seconds"),
            "output_dir": str(root),
            "ok": False,
            "issues": issues,
            "details": details,
        }

    manifest_path = root / "generation_manifest.json"
    if not manifest_path.exists():
        add_issue(issues, "error", "manifest", "缺少 generation_manifest.json。")
        manifest = {}
    else:
        manifest = read_json(manifest_path)
        details["manifest"] = {
            "topic": (manifest.get("request") or {}).get("topic", ""),
            "grade": (manifest.get("request") or {}).get("grade", ""),
            "outputs": (manifest.get("request") or {}).get("outputs", []),
            "skip_media": (manifest.get("generation_options") or {}).get("skip_media", False),
            "include_video": (manifest.get("generation_options") or {}).get("include_video", False),
        }

    outputs = manifest.get("outputs", {}) if manifest else {}
    for key in REQUIRED_OUTPUT_KEYS:
        value = outputs.get(key)
        if not value:
            add_issue(issues, "error", key, f"manifest 缺少輸出項目：{key}")
            continue
        path = resolve_output_path(root, value)
        if not path.exists():
            add_issue(issues, "error", key, f"輸出檔不存在：{path}")
            continue
        if path.stat().st_size <= 0:
            add_issue(issues, "error", key, f"輸出檔是空檔：{path}")
            continue
        details[key] = {"path": str(path), "bytes": path.stat().st_size}

    lesson_path = root / "lesson_structure.json"
    lesson = {}
    if lesson_path.exists():
        try:
            lesson = read_json(lesson_path)
            details["lesson_structure"] = {
                **details.get("lesson_structure", {}),
                "title": lesson.get("title", ""),
                "vocabulary_count": len(lesson.get("vocabulary", []) or []),
                "dialogue_count": len(lesson.get("dialogues", []) or []),
                "question_count": len(lesson.get("questions", []) or []),
                "has_official_recommendations": bool(lesson.get("official_material_recommendations")),
            }
            if not lesson.get("questions"):
                add_issue(issues, "error", "lesson_structure", "教材結構缺少測驗題。")
            if not lesson.get("vocabulary"):
                add_issue(issues, "warning", "lesson_structure", "教材結構缺少詞彙。")
            details["topic_alignment"] = validate_topic_alignment(lesson, manifest, issues)
        except Exception as exc:
            add_issue(issues, "error", "lesson_structure", f"lesson_structure.json 無法解析：{exc}")

    for key in ("student_worksheet", "teacher_guide", "exam_paper", "exam_answer_key"):
        path_text = outputs.get(key)
        if path_text:
            path = resolve_output_path(root, path_text)
            if path.exists():
                details[key] = {**details.get(key, {}), **validate_docx(path, issues, key)}

    slides_path = outputs.get("slides")
    if slides_path:
        path = resolve_output_path(root, slides_path)
        if path.exists():
            details["slides"] = {**details.get("slides", {}), **validate_pptx(path, issues)}

    html_path = outputs.get("interactive_website")
    if html_path:
        path = resolve_output_path(root, html_path)
        if path.exists():
            details["interactive_website"] = {
                **details.get("interactive_website", {}),
                **validate_html(path, issues),
            }

    requested_outputs = (manifest.get("request") or {}).get("outputs", []) if manifest else []
    if "quiz" in requested_outputs:
        quiz_path_text = outputs.get("quiz_bank")
        if not quiz_path_text:
            add_issue(issues, "error", "quiz_bank", "自然語言需求包含測驗，但 manifest 缺少 quiz_bank。")
        else:
            quiz_path = resolve_output_path(root, quiz_path_text)
            if not quiz_path.exists():
                add_issue(issues, "error", "quiz_bank", f"測驗題庫不存在：{quiz_path}")
            elif quiz_path.stat().st_size <= 0:
                add_issue(issues, "error", "quiz_bank", f"測驗題庫是空檔：{quiz_path}")
            else:
                details["quiz_bank"] = {
                    "path": str(quiz_path),
                    "bytes": quiz_path.stat().st_size,
                    **validate_quiz_bank(quiz_path, issues),
                }

        key_path_text = outputs.get("quiz_teacher_key")
        if not key_path_text:
            add_issue(issues, "error", "quiz_teacher_key", "自然語言需求包含測驗，但 manifest 缺少 quiz_teacher_key。")
        else:
            key_path = resolve_output_path(root, key_path_text)
            if not key_path.exists():
                add_issue(issues, "error", "quiz_teacher_key", f"測驗教師答案不存在：{key_path}")
            elif key_path.stat().st_size <= 0:
                add_issue(issues, "error", "quiz_teacher_key", f"測驗教師答案是空檔：{key_path}")
            else:
                details["quiz_teacher_key"] = {"path": str(key_path), "bytes": key_path.stat().st_size}

    video_generation = manifest.get("video_generation", {}) if manifest else {}
    if video_generation:
        details["video_generation"] = video_generation
        video_path_text = video_generation.get("path") or outputs.get("video")
        if video_generation.get("attempted"):
            if not video_generation.get("success"):
                add_issue(issues, "error", "video", "影片已嘗試產生但失敗。")
            elif not video_path_text:
                add_issue(issues, "error", "video", "影片產生標示成功，但 manifest 沒有影片路徑。")
            else:
                video_path = resolve_output_path(root, video_path_text)
                if not video_path.exists():
                    add_issue(issues, "error", "video", f"影片檔不存在：{video_path}")
                elif video_path.stat().st_size <= 0:
                    add_issue(issues, "error", "video", f"影片檔是空檔：{video_path}")
                else:
                    details["video"] = {"path": str(video_path), "bytes": video_path.stat().st_size}
        elif video_generation.get("requested_in_text") and video_generation.get("skipped_reason") == "disabled_by_option":
            add_issue(issues, "warning", "video", "自然語言需求包含影片，但本次已用選項跳過影片生成。")
    elif "video" in ((manifest.get("request") or {}).get("outputs", []) if manifest else []):
        video_path_text = outputs.get("video")
        if video_path_text:
            video_path = resolve_output_path(root, video_path_text)
            if not video_path.exists() or video_path.stat().st_size <= 0:
                add_issue(issues, "error", "video", f"自然語言需求包含影片，但影片檔不存在或為空：{video_path}")
        else:
            add_issue(issues, "warning", "video", "自然語言需求包含影片，但 manifest 未記錄影片生成狀態。")

    recommendations = manifest.get("official_material_recommendations", {}) if manifest else {}
    if not recommendations:
        add_issue(issues, "warning", "official_material_recommendations", "manifest 沒有官方素材推薦。")
    else:
        for output in requested_outputs:
            if output != "video" and not recommendations.get(output):
                add_issue(
                    issues,
                    "warning",
                    "official_material_recommendations",
                    f"缺少「{output}」的官方素材推薦。"
                )

    ok = not any(issue["severity"] == "error" for issue in issues)
    return {
        "validated_at": datetime.now().isoformat(timespec="seconds"),
        "output_dir": str(root),
        "ok": ok,
        "issues": issues,
        "details": details,
    }


def write_reports(result: Dict[str, Any], output_dir: Path) -> None:
    json_path = output_dir / "generation_validation.json"
    md_path = output_dir / "generation_validation.md"
    json_path.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")

    lines = [
        "# 教材產出驗證報告",
        "",
        f"驗證時間：{result['validated_at']}",
        f"資料夾：`{result['output_dir']}`",
        f"結果：{'通過' if result['ok'] else '未通過'}",
        "",
        "## 問題",
        "",
    ]
    if result["issues"]:
        for issue in result["issues"]:
            lines.append(f"- [{issue['severity']}] {issue['item']}：{issue['message']}")
    else:
        lines.append("- 無")
    lines.extend(["", "## 摘要", ""])
    manifest = result.get("details", {}).get("manifest", {})
    if manifest:
        lines.append(f"- 主題：{manifest.get('topic', '')}")
        lines.append(f"- 年級：{manifest.get('grade', '')}")
        lines.append(f"- 產出類型：{', '.join(manifest.get('outputs', []))}")
        lines.append(f"- 快速模式：{manifest.get('skip_media', False)}")
    lesson = result.get("details", {}).get("lesson_structure", {})
    if lesson:
        lines.append(f"- 詞彙數：{lesson.get('vocabulary_count', 0)}")
        lines.append(f"- 對話數：{lesson.get('dialogue_count', 0)}")
        lines.append(f"- 題目數：{lesson.get('question_count', 0)}")
    alignment = result.get("details", {}).get("topic_alignment", {})
    if alignment:
        matched = ", ".join(alignment.get("matched_terms", [])) or "無"
        lines.append(f"- 主題貼合：{'通過' if alignment.get('is_aligned') else '未通過'}（命中：{matched}）")
    slides = result.get("details", {}).get("slides", {})
    if slides:
        lines.append(f"- 簡報投影片數：{slides.get('slide_count', 0)}")
    quiz_bank = result.get("details", {}).get("quiz_bank", {})
    if quiz_bank:
        lines.append(f"- 測驗題庫：{quiz_bank.get('question_count', 0)} 題自動計分，{quiz_bank.get('official_extension_count', 0)} 題官方延伸")
    video_generation = result.get("details", {}).get("video_generation", {})
    if video_generation:
        status = "成功" if video_generation.get("success") else ("已嘗試但失敗" if video_generation.get("attempted") else "未產生")
        if video_generation.get("skipped_reason") == "disabled_by_option":
            status = "已用選項跳過"
        lines.append(f"- 影片生成：{status}")
    md_path.write_text("\n".join(lines), encoding="utf-8")


def main() -> int:
    if len(sys.argv) < 2:
        print("用法：python validate_generation_output.py <輸出資料夾>")
        return 2
    output_dir = Path(sys.argv[1])
    result = validate_output_folder(str(output_dir))
    if output_dir.exists():
        write_reports(result, output_dir)
    print(json.dumps({"ok": result["ok"], "issues": result["issues"]}, ensure_ascii=False, indent=2))
    return 0 if result["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
