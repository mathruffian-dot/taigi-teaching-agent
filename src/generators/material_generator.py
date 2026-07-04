# 臺語教材生成器主程式 (material_generator.py)
import os
import sys
import json
import argparse
from datetime import datetime
from typing import Dict, Any

# 加入專案 src 目錄到路徑
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from utils import safe_print, resolve_output_base_dir
print = safe_print

from rag.retriever import TaigiRetriever
from tailo.validator import convert_sentence_numeric_to_diacritic
from tailo.piauim import Piauim
from tts.generator import TaigiTTS
from generators.image_generator import FreeImageGenerator
from agent.content_checker import check_lesson_content

class MaterialGenerator:
    def __init__(self, config_path: str = "config.json"):
        self.config_path = config_path
        self.config = self._load_config(config_path)
        self.retriever = TaigiRetriever()
        self.tts = TaigiTTS(config_path)
        self.piauim = Piauim(self.config)

    def _load_config(self, filepath: str) -> Dict[str, Any]:
        if os.path.exists(filepath):
            with open(filepath, "r", encoding="utf-8-sig") as f:
                return json.load(f)
        return {}

    def _format_official_material_line(self, item: Dict[str, Any]) -> str:
        title = item.get("title", "")
        stage = item.get("learning_stage", "")
        kind = item.get("resource_kind") or item.get("material_type", "")
        provider = item.get("provider", "")
        parts = [part for part in [title, stage, kind, provider] if part]
        return "｜".join(parts)

    def _flatten_recommended_materials(self, data: Dict[str, Any], limit: int = 6):
        grouped = data.get("official_material_recommendations", {}) or {}
        if not grouped:
            return data.get("official_materials", [])[:limit]

        selected = []
        seen = set()
        for output in ("exam", "worksheet", "slides", "video", "interactive", "quiz"):
            for item in grouped.get(output, []) or []:
                key = item.get("page_url") or item.get("title")
                if key in seen:
                    continue
                seen.add(key)
                selected.append(item)
                if len(selected) >= limit:
                    return selected
        return selected

    def _flatten_generation_assets(self, data: Dict[str, Any], output: str = None, asset_type: str = None, limit: int = 6):
        grouped = data.get("official_generation_assets", {}) or {}
        selected = []
        seen = set()
        groups = [output] if output else ["exam", "worksheet", "slides", "video", "interactive", "quiz"]
        for group in groups:
            for asset in grouped.get(group, []) or []:
                if asset_type and asset.get("asset_type") != asset_type:
                    continue
                key = asset.get("asset_id") or f"{asset.get('asset_type')}|{asset.get('title')}|{asset.get('question')}"
                if key in seen:
                    continue
                seen.add(key)
                selected.append(asset)
                if len(selected) >= limit:
                    return selected
        return selected

    def _attach_official_generation_assets(self, data: Dict[str, Any]) -> None:
        data["official_generated_questions"] = self._flatten_generation_assets(
            data, output="exam", asset_type="multiple_choice", limit=5
        )
        data["official_slide_seeds"] = self._flatten_generation_assets(
            data, output="slides", asset_type="slide_seed", limit=4
        )
        data["official_interactive_assets"] = self._flatten_generation_assets(
            data, output="interactive", limit=5
        )

    def generate_all(self, case_path: str, output_dir: str = None, skip_media: bool = False):
        # 輸出目錄優先序：呼叫參數 > config 的 output.base_dir > 預設 "output"
        # 大型產出（音訊/影片）建議透過 config 指向本機目錄，避免雲端硬碟同步失敗。
        if output_dir is None:
            output_dir = resolve_output_base_dir(self.config)
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        print(f"[*] 輸出目錄: {output_dir}")

        # 1. 讀取輸入測試案例
        with open(case_path, "r", encoding="utf-8") as f:
            case_data = json.load(f)

        print(f"[*] 讀取測試案例: {case_data.get('title', '未命名')}")
        if skip_media:
            print("[*] 快速模式：跳過音訊與圖片產生，只輸出文字、考卷、簡報與互動網頁。")

        # 2. 透過 RAG 豐富課綱與詞彙庫資料
        print("[*] 執行 RAG 知識庫檢索與欄位豐富化...")
        enriched_data = self.retriever.enrich_lesson_json(case_data)
        if not enriched_data.get("official_materials"):
            official_query = enriched_data.get("title") or case_data.get("title", "")
            enriched_data["official_materials"] = self.retriever.retrieve_official_materials(official_query, limit=5)
        if not enriched_data.get("official_material_recommendations"):
            request = enriched_data.get("natural_language_request", {}) or {}
            official_query = request.get("topic") or enriched_data.get("title") or case_data.get("title", "")
            outputs = request.get("outputs") or ["worksheet", "slides", "interactive", "quiz"]
            enriched_data["official_material_recommendations"] = self.retriever.recommend_official_materials(
                official_query,
                outputs=outputs,
                limit_per_output=3,
            )
        if not enriched_data.get("official_generation_assets") and hasattr(self.retriever, "recommend_official_generation_assets"):
            request = enriched_data.get("natural_language_request", {}) or {}
            official_query = request.get("topic") or enriched_data.get("title") or case_data.get("title", "")
            outputs = request.get("outputs") or ["worksheet", "slides", "interactive", "quiz"]
            enriched_data["official_generation_assets"] = self.retriever.recommend_official_generation_assets(
                official_query,
                outputs=outputs,
                limit_per_output=3,
            )
        self._attach_official_generation_assets(enriched_data)

        # 3. 產生台羅調符式。優先用意傳「漢字→台羅」權威標音（取代 LLM 可能不準的拼音），
        #    意傳查不到才退回把 LLM 的數字調轉成調符式。
        use_piauim = self.config.get("piauim", {}).get("provider", "ithuan") != "off"
        if use_piauim:
            print("[*] 以意傳標音由漢字重標台羅（取代 LLM 拼音）...")
        else:
            print("[*] 處理台羅拼音數值調至調符式之自動轉換...")

        def _standardize_tailo(item):
            hanji = item.get("hanji", "")
            kip = self.piauim.kip(hanji) if (use_piauim and hanji) else None
            if kip:
                item["tailo_diacritic"] = kip
                item["tailo_source"] = "ithuan"
            elif item.get("tailo_numeric") and item.get("tailo_numeric") != "pending":
                item["tailo_diacritic"] = convert_sentence_numeric_to_diacritic(item["tailo_numeric"])
                item["tailo_source"] = "llm"

        for vocab in enriched_data.get("vocabulary", []):
            _standardize_tailo(vocab)
        for dia in enriched_data.get("dialogues", []):
            _standardize_tailo(dia)

        # 3.1. 內容自動檢核（標記華語用字、漢字/臺羅音節不符等，供教師審核）
        content_warnings = check_lesson_content(enriched_data)
        if content_warnings:
            print(f"[!] 內容檢核發現 {len(content_warnings)} 項待確認（已寫入教師審核報告）：")
            for w in content_warnings:
                print(f"    ⚠️ {w}")
        else:
            print("[*] 內容檢核通過，未發現明顯用字/拼音問題。")

        if skip_media:
            for vocab in enriched_data.get("vocabulary", []):
                vocab["audio_file"] = ""
                vocab["image_file"] = ""
            for dia in enriched_data.get("dialogues", []):
                dia["audio_file"] = ""
        else:
            # 3.5. 產生語音音訊檔 (詞彙下載與對話合成)
            print("[*] 執行語音生成與下載...")
            audio_dir = os.path.join(output_dir, "audio")
            if not os.path.exists(audio_dir):
                os.makedirs(audio_dir)
                
            if self.tts.provider == "voxcpm":
                # 全教材統一用「三師爸台語」聲音（VoxCPM2 批次，模型只載一次）。
                # 詞彙不再抓萌典官方單詞音檔，改由三師爸發音，全份語氣一致。
                batch = []
                refs = []  # [(item, audio_relpath)]
                for vocab in enriched_data.get("vocabulary", []):
                    hanji = vocab.get("hanji", "")
                    if not hanji:
                        continue
                    wav_filename = f"vocab_{hanji}.wav"
                    feed = self.tts._voxcpm_feed_text(hanji, vocab.get("tailo_diacritic", ""))
                    batch.append({"text": feed, "output": os.path.join(audio_dir, wav_filename)})
                    refs.append((vocab, f"audio/{wav_filename}"))
                for idx, dia in enumerate(enriched_data.get("dialogues", [])):
                    hanji = dia.get("hanji", "")
                    if not hanji:
                        continue
                    wav_filename = f"dialogue_{idx}.wav"
                    feed = self.tts._voxcpm_feed_text(hanji, dia.get("tailo_diacritic", ""))
                    batch.append({"text": feed, "output": os.path.join(audio_dir, wav_filename)})
                    refs.append((dia, f"audio/{wav_filename}"))

                if batch:
                    flags = self.tts.synthesize_voxcpm_batch(batch)
                    for (item, rel), ok in zip(refs, flags):
                        item["audio_file"] = rel  # 即便個別降級，仍指向該檔
            else:
                for vocab in enriched_data.get("vocabulary", []):
                    hanji = vocab.get("hanji", "")
                    if hanji:
                        # 預設儲存為 ogg，若失敗則合成 wav 佔位符
                        ogg_filename = f"vocab_{hanji}.ogg"
                        ogg_path = os.path.join(audio_dir, ogg_filename)

                        # 嘗試從萌典下載
                        success = self.tts.fetch_vocab_audio(hanji, ogg_path)
                        if success:
                            vocab["audio_file"] = f"audio/{ogg_filename}"
                        else:
                            # 降級：使用 dummy/yating 合成 wav 檔
                            wav_filename = f"vocab_{hanji}.wav"
                            wav_path = os.path.join(audio_dir, wav_filename)
                            self.tts.synthesize_sentence(hanji, wav_path, vocab.get("tailo_numeric", ""),
                                                         vocab.get("tailo_diacritic", ""))
                            vocab["audio_file"] = f"audio/{wav_filename}"

                for idx, dia in enumerate(enriched_data.get("dialogues", [])):
                    hanji = dia.get("hanji", "")
                    if hanji:
                        # 對話句子進行語音合成
                        wav_filename = f"dialogue_{idx}.wav"
                        wav_path = os.path.join(audio_dir, wav_filename)
                        self.tts.synthesize_sentence(hanji, wav_path, dia.get("tailo_numeric", ""),
                                                     dia.get("tailo_diacritic", ""))
                        dia["audio_file"] = f"audio/{wav_filename}"

            # 3.6. 產生詞彙插圖 (免費生圖 API Option B)
            print("[*] 執行免費 AI 詞彙生圖...")
            image_dir = os.path.join(output_dir, "images")
            if not os.path.exists(image_dir):
                os.makedirs(image_dir)

            img_gen = FreeImageGenerator(self.config_path)

            for vocab in enriched_data.get("vocabulary", []):
                hanji = vocab.get("hanji", "")
                zh_tw = vocab.get("zh_tw", "")
                if not zh_tw or zh_tw == "pending":
                    zh_tw = hanji

                if zh_tw:
                    img_filename = f"vocab_{hanji}.jpg"
                    img_path = os.path.join(image_dir, img_filename)
                    success = img_gen.generate_image(zh_tw, img_path)
                    if success:
                        vocab["image_file"] = f"images/{img_filename}"
                    else:
                        vocab["image_file"] = ""

        enriched_data["generation_options"] = {
            **(enriched_data.get("generation_options", {}) or {}),
            "skip_media": skip_media,
        }

        # 保存豐富化後的教材結構資料
        json_output = os.path.join(output_dir, "lesson_structure.json")
        with open(json_output, "w", encoding="utf-8") as f:
            json.dump(enriched_data, f, ensure_ascii=False, indent=2)
        print(f"  [+] 已產生教材結構 JSON: {json_output}")

        # 4. 產生講義 (Word DOCX)
        self._generate_docx(enriched_data, output_dir)

        # 5. 產生考卷 (Word DOCX)
        self._generate_exam_docx(enriched_data, output_dir)

        # 6. 產生獨立測驗題庫 (JSON + 教師答案 Markdown)
        self._generate_quiz_bank(enriched_data, output_dir)

        # 7. 產生簡報 (PowerPoint PPTX)
        self._generate_pptx(enriched_data, output_dir)

        # 8. 產生互動網站 (離線單一 HTML)
        self._generate_html(enriched_data, output_dir)

        # 9. 產生教師審核報告
        self._generate_review_report(enriched_data, output_dir)
        return output_dir

    def _generate_docx(self, data: Dict[str, Any], output_dir: str):
        try:
            from docx import Document
        except ImportError:
            print("  [-] 警告: 找不到 python-docx 模組，跳過 Word 檔案生成。")
            return
            
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls
            
        def set_cell_shading(cell, color_hex):
            shading_xml = f'<w:shd {nsdecls("w")} w:fill="{color_hex}"/>'
            cell._tc.get_or_add_tcPr().append(parse_xml(shading_xml))
            
        def style_run(run, font_name="Arial", size_pt=11, bold=False, color_rgb=None):
            run.font.name = font_name
            run.font.size = Pt(size_pt)
            run.bold = bold
            if color_rgb:
                run.font.color.rgb = color_rgb

        # ==================== 學生版講義 ====================
        doc_student = Document()
        title_p = doc_student.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title_p.add_run(f"臺語學習講義：{data.get('title')}")
        style_run(run, "微軟正黑體", 18, True, RGBColor(11, 60, 48))
        
        meta_p = doc_student.add_paragraph()
        meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = meta_p.add_run(f"適用年級：{data.get('grade')} | 課程時間：{data.get('duration_minutes')} 分鐘 | 姓名：___________ 座號：_____")
        style_run(run, "微軟正黑體", 10, False, RGBColor(82, 100, 95))
        
        # 1. 課綱指標
        h1 = doc_student.add_paragraph()
        style_run(h1.add_run("一、學習目標與課綱指標對照"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
        for perf in data.get("curriculum", {}).get("learning_performance", []):
            p = doc_student.add_paragraph()
            style_run(p.add_run(f"• 學習表現：{perf}"), "新細明體", 10, False)
        for cont in data.get("curriculum", {}).get("learning_content", []):
            p = doc_student.add_paragraph()
            style_run(p.add_run(f"• 學習內容：{cont}"), "新細明體", 10, False)
            
        # 2. 詞彙表 (美化斑馬紋表格)
        h2 = doc_student.add_paragraph()
        style_run(h2.add_run("二、核心詞彙認讀與手寫練習"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
        
        table = doc_student.add_table(rows=1, cols=5)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = '圖示'
        hdr_cells[1].text = '臺語漢字'
        hdr_cells[2].text = '教育部臺羅拼音'
        hdr_cells[3].text = '華語翻譯'
        hdr_cells[4].text = '手寫練習 (漢字與拼音)'
        
        # 設定表頭顏色 (主題墨綠)
        for cell in hdr_cells:
            set_cell_shading(cell, "0B3C30")
            for p in cell.paragraphs:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                for r in p.runs:
                    style_run(r, "微軟正黑體", 10, True, RGBColor(255, 255, 255))
        
        row_count = 0
        for vocab in data.get("vocabulary", []):
            row_cells = table.add_row().cells
            
            # 1. 置入圖片
            from docx.shared import Inches
            img_rel_path = vocab.get("image_file")
            if img_rel_path:
                full_img_path = os.path.join(output_dir, img_rel_path)
                if os.path.exists(full_img_path):
                    p_img = row_cells[0].paragraphs[0]
                    p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run_img = p_img.add_run()
                    try:
                        run_img.add_picture(full_img_path, width=Inches(0.9))
                    except Exception as e:
                        p_img.text = "[圖片損毀]"
                else:
                    row_cells[0].text = "無圖片"
            else:
                row_cells[0].text = "無圖片"
                
            row_cells[1].text = vocab.get("hanji", "")
            row_cells[2].text = vocab.get("tailo_diacritic", "")
            row_cells[3].text = vocab.get("zh_tw", "")
            row_cells[4].text = "__________________"
            
            # 斑馬紋底色
            row_color = "F6F8F6" if row_count % 2 == 1 else "FFFFFF"
            for i, cell in enumerate(row_cells):
                set_cell_shading(cell, row_color)
                p = cell.paragraphs[0]
                # 前四欄置中，手寫欄靠左
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER if i < 4 else WD_ALIGN_PARAGRAPH.LEFT
                for r in p.runs:
                    style_run(r, "新細明體" if i != 2 else "Arial", 10, False, RGBColor(47, 62, 70))
            row_count += 1
            
        # 3. 情境對話
        doc_student.add_paragraph() # 空行
        h3 = doc_student.add_paragraph()
        style_run(h3.add_run("三、情境會話認讀 (口說與聆聽)"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
        for dia in data.get("dialogues", []):
            p = doc_student.add_paragraph()
            r_role = p.add_run(f"🗣️ {dia.get('role')}：")
            style_run(r_role, "微軟正黑體", 11, True, RGBColor(184, 134, 11))
            r_text = p.add_run(dia.get('hanji'))
            style_run(r_text, "新細明體", 11, False)
            
            p_py = doc_student.add_paragraph()
            r_py = p_py.add_run(f"   [{dia.get('tailo_diacritic')}]")
            style_run(r_py, "Arial", 10, False, RGBColor(82, 100, 95))
            
            p_zh = doc_student.add_paragraph()
            r_zh = p_zh.add_run(f"   (請寫出此句華語意譯：__________________________________)")
            style_run(r_zh, "新細明體", 10, False, RGBColor(127, 140, 141))
            
        # 4. 評量題目
        doc_student.add_paragraph()
        h4 = doc_student.add_paragraph()
        style_run(h4.add_run("四、課堂自我檢測 (隨堂測驗)"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
        q_count = 1
        for q in data.get("questions", []):
            p = doc_student.add_paragraph()
            style_run(p.add_run(f"({q_count}) {q.get('question')}"), "微軟正黑體", 11, True)
            for idx, opt in enumerate(q.get("options", [])):
                p_opt = doc_student.add_paragraph()
                style_run(p_opt.add_run(f"    [  ] {opt}"), "新細明體", 10, False)
            q_count += 1

        official_questions = data.get("official_generated_questions", []) or []
        if official_questions:
            doc_student.add_paragraph()
            h5 = doc_student.add_paragraph()
            style_run(h5.add_run("五、官方教材延伸題（教師確認答案後使用）"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
            for asset in official_questions:
                p = doc_student.add_paragraph()
                style_run(p.add_run(f"({q_count}) {asset.get('question', '')}"), "微軟正黑體", 11, True)
                for opt in asset.get("options", []) or []:
                    p_opt = doc_student.add_paragraph()
                    style_run(p_opt.add_run(f"    [  ] {opt}"), "新細明體", 10, False)
                source_p = doc_student.add_paragraph()
                style_run(source_p.add_run(f"    來源：{asset.get('title', '')}；答案需教師確認。"), "新細明體", 9, False, RGBColor(127, 140, 141))
                q_count += 1
            
        student_path = os.path.join(output_dir, "student_worksheet.docx")
        doc_student.save(student_path)
        print(f"  [+] 已產生學生版講義 Word: {student_path}")
        
        # ==================== 教師解答版講義 ====================
        doc_teacher = Document()
        title_p = doc_teacher.add_paragraph()
        title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = title_p.add_run(f"臺語學習講義 (教師解答指導版)：{data.get('title')}")
        style_run(run, "微軟正黑體", 18, True, RGBColor(184, 134, 11))
        
        h1 = doc_teacher.add_paragraph()
        style_run(h1.add_run("一、核心對話與翻譯解答"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
        for dia in data.get("dialogues", []):
            p = doc_teacher.add_paragraph()
            style_run(p.add_run(f"🗣️ {dia.get('role')}：{dia.get('hanji')}"), "微軟正黑體", 11, True)
            p_ans = doc_teacher.add_paragraph()
            r_ans = p_ans.add_run(f"   [翻譯解答] {dia.get('zh_tw')}")
            style_run(r_ans, "微軟正黑體", 10, True, RGBColor(180, 40, 40))
            
        h2 = doc_teacher.add_paragraph()
        style_run(h2.add_run("二、評量測驗答案與解析"), "微軟正黑體", 14, True, RGBColor(11, 60, 48))
        q_count = 1
        for q in data.get("questions", []):
            p = doc_teacher.add_paragraph()
            style_run(p.add_run(f"({q_count}) {q.get('question')}"), "微軟正黑體", 11, True)
            
            ans_text = q.get("options")[q.get("answer_index")]
            p_ans = doc_teacher.add_paragraph()
            r_ans = p_ans.add_run(f"   ★ 正確答案: {ans_text}")
            style_run(r_ans, "微軟正黑體", 10, True, RGBColor(180, 40, 40))
            
            p_exp = doc_teacher.add_paragraph()
            r_exp = p_exp.add_run(f"   [解析說明] {q.get('explanation')}")
            style_run(r_exp, "新細明體", 10, False, RGBColor(120, 120, 120))
            q_count += 1
            
        teacher_path = os.path.join(output_dir, "teacher_guide.docx")
        doc_teacher.save(teacher_path)
        print(f"  [+] 已產生教師版講義 Word: {teacher_path}")

    def _generate_exam_docx(self, data: Dict[str, Any], output_dir: str):
        try:
            from docx import Document
        except ImportError:
            print("  [-] 警告: 找不到 python-docx 模組，跳過考卷生成。")
            return

        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        def style_run(run, font_name="微軟正黑體", size_pt=11, bold=False, color_rgb=None):
            run.font.name = font_name
            run.font.size = Pt(size_pt)
            run.bold = bold
            if color_rgb:
                run.font.color.rgb = color_rgb

        def add_header(doc, title):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            style_run(p.add_run(title), "微軟正黑體", 18, True, RGBColor(11, 60, 48))
            meta = doc.add_paragraph()
            meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
            style_run(
                meta.add_run(f"適用年級：{data.get('grade', '')}｜姓名：__________｜座號：_____｜分數：_____"),
                "微軟正黑體",
                10,
                False,
                RGBColor(82, 100, 95),
            )

        exam = Document()
        add_header(exam, f"臺語單元考卷：{data.get('title', '')}")
        exam.add_paragraph("一、選擇題").runs[0].bold = True
        for idx, question in enumerate(data.get("questions", []), start=1):
            p = exam.add_paragraph()
            style_run(p.add_run(f"{idx}. {question.get('question', '')}"), "微軟正黑體", 11, True)
            for opt_idx, option in enumerate(question.get("options", [])):
                label = chr(ord("A") + opt_idx)
                opt_p = exam.add_paragraph()
                style_run(opt_p.add_run(f"   ({label}) {option}"), "新細明體", 10, False)

        official_questions = data.get("official_generated_questions", []) or []
        if official_questions:
            exam.add_paragraph("")
            exam.add_paragraph("官方教材延伸題（教師確認答案後使用）").runs[0].bold = True
            start_idx = len(data.get("questions", [])) + 1
            for idx, asset in enumerate(official_questions, start=start_idx):
                p = exam.add_paragraph()
                style_run(p.add_run(f"{idx}. {asset.get('question', '')}"), "微軟正黑體", 11, True)
                for opt_idx, option in enumerate(asset.get("options", []) or []):
                    label = chr(ord("A") + opt_idx)
                    opt_p = exam.add_paragraph()
                    style_run(opt_p.add_run(f"   ({label}) {option}"), "新細明體", 10, False)
                source_p = exam.add_paragraph()
                style_run(source_p.add_run(f"   來源：{asset.get('title', '')}；答案需教師確認。"), "新細明體", 9, False, RGBColor(127, 140, 141))

        exam.add_paragraph("")
        exam.add_paragraph("二、詞彙書寫").runs[0].bold = True
        for idx, vocab in enumerate(data.get("vocabulary", [])[:5], start=1):
            p = exam.add_paragraph()
            style_run(
                p.add_run(f"{idx}. {vocab.get('zh_tw', '')}：臺語漢字 ____________，臺羅 ____________"),
                "新細明體",
                10,
                False,
            )

        exam_path = os.path.join(output_dir, "exam_paper.docx")
        exam.save(exam_path)
        print(f"  [+] 已產生學生考卷 Word: {exam_path}")

        answer_key = Document()
        add_header(answer_key, f"臺語單元考卷解答：{data.get('title', '')}")
        answer_key.add_paragraph("一、選擇題解答").runs[0].bold = True
        for idx, question in enumerate(data.get("questions", []), start=1):
            options = question.get("options", []) or []
            answer_index = question.get("answer_index", 0)
            answer = options[answer_index] if 0 <= answer_index < len(options) else ""
            p = answer_key.add_paragraph()
            style_run(p.add_run(f"{idx}. {answer}"), "微軟正黑體", 11, True, RGBColor(180, 40, 40))
            exp = question.get("explanation", "")
            if exp:
                style_run(answer_key.add_paragraph().add_run(f"   解析：{exp}"), "新細明體", 10, False)

        official_questions = data.get("official_generated_questions", []) or []
        if official_questions:
            answer_key.add_paragraph("")
            answer_key.add_paragraph("官方教材延伸題").runs[0].bold = True
            start_idx = len(data.get("questions", [])) + 1
            for idx, asset in enumerate(official_questions, start=start_idx):
                p = answer_key.add_paragraph()
                style_run(p.add_run(f"{idx}. {asset.get('question', '')}"), "微軟正黑體", 10, True, RGBColor(184, 134, 11))
                style_run(answer_key.add_paragraph().add_run("   答案：需教師依官方教材確認。"), "新細明體", 10, False)

        answer_key.add_paragraph("")
        answer_key.add_paragraph("二、詞彙書寫參考答案").runs[0].bold = True
        for idx, vocab in enumerate(data.get("vocabulary", [])[:5], start=1):
            p = answer_key.add_paragraph()
            style_run(
                p.add_run(f"{idx}. {vocab.get('zh_tw', '')}：{vocab.get('hanji', '')}｜{vocab.get('tailo_diacritic') or vocab.get('tailo_numeric', '')}"),
                "新細明體",
                10,
                False,
            )

        answer_path = os.path.join(output_dir, "exam_answer_key.docx")
        answer_key.save(answer_path)
        print(f"  [+] 已產生考卷解答 Word: {answer_path}")

    def _generate_quiz_bank(self, data: Dict[str, Any], output_dir: str):
        questions = []
        for idx, question in enumerate(data.get("questions", []) or [], start=1):
            options = question.get("options", []) or []
            answer_index = question.get("answer_index", 0)
            answer = options[answer_index] if 0 <= answer_index < len(options) else ""
            questions.append({
                "id": question.get("id") or f"q{idx}",
                "type": "multiple_choice",
                "question": question.get("question", ""),
                "options": options,
                "answer_index": answer_index,
                "answer": answer,
                "explanation": question.get("explanation", ""),
                "graded": True,
                "source": "generated_lesson",
            })

        official_extension_questions = []
        for idx, asset in enumerate(data.get("official_generated_questions", []) or [], start=1):
            official_extension_questions.append({
                "id": asset.get("asset_id") or f"official_q{idx}",
                "type": asset.get("asset_type") or "multiple_choice",
                "question": asset.get("question", ""),
                "options": asset.get("options", []) or [],
                "answer_index": None,
                "answer": None,
                "explanation": "需教師依官方教材確認正確答案後使用。",
                "graded": False,
                "teacher_answer_required": True,
                "source": "official_generation_asset",
                "source_title": asset.get("title", ""),
                "source_url": asset.get("source_url") or asset.get("page_url", ""),
            })

        quiz_bank = {
            "title": data.get("title", ""),
            "grade": data.get("grade", ""),
            "duration_minutes": data.get("duration_minutes", ""),
            "generated_at": datetime.now().isoformat(timespec="seconds"),
            "question_count": len(questions),
            "official_extension_count": len(official_extension_questions),
            "questions": questions,
            "official_extension_questions": official_extension_questions,
        }

        quiz_path = os.path.join(output_dir, "quiz_bank.json")
        with open(quiz_path, "w", encoding="utf-8") as f:
            json.dump(quiz_bank, f, ensure_ascii=False, indent=2)

        key_lines = [
            f"# 測驗題庫教師答案：{data.get('title', '')}",
            "",
            f"- 適用年級：{data.get('grade', '')}",
            f"- 自動計分題數：{len(questions)}",
            f"- 官方延伸題數：{len(official_extension_questions)}",
            "",
            "## 自動計分題",
            "",
        ]
        for idx, question in enumerate(questions, start=1):
            key_lines.append(f"{idx}. {question['question']}")
            key_lines.append(f"   - 答案：{question['answer']}")
            if question.get("explanation"):
                key_lines.append(f"   - 解析：{question['explanation']}")
        if official_extension_questions:
            key_lines.extend(["", "## 官方延伸題（教師確認後使用）", ""])
            for idx, question in enumerate(official_extension_questions, start=1):
                key_lines.append(f"{idx}. {question['question']}")
                key_lines.append("   - 答案：需教師依官方教材確認。")
                if question.get("source_title"):
                    key_lines.append(f"   - 來源：{question['source_title']}")

        key_path = os.path.join(output_dir, "quiz_teacher_key.md")
        with open(key_path, "w", encoding="utf-8") as f:
            f.write("\n".join(key_lines))

        print(f"  [+] 已產生測驗題庫 JSON: {quiz_path}")
        print(f"  [+] 已產生測驗教師答案 Markdown: {key_path}")

    def _generate_pptx(self, data: Dict[str, Any], output_dir: str):
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
        except ImportError:
            print("  [-] 警告: 找不到 python-pptx 模組，跳過 PowerPoint 簡報生成。")
            return

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        theme = {
            "green": RGBColor(11, 60, 48),
            "gold": RGBColor(184, 134, 11),
            "cream": RGBColor(250, 248, 239),
            "ink": RGBColor(47, 62, 70),
            "muted": RGBColor(95, 112, 105),
        }

        def set_bg(slide, color):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_textbox(slide, left, top, width, height, text="", size=24, bold=False, color=None):
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            frame = box.text_frame
            frame.clear()
            p = frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.name = "Microsoft JhengHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or theme["ink"]
            return box

        def add_title(slide, title, subtitle=""):
            add_textbox(slide, 0.65, 0.35, 11.8, 0.55, title, 24, True, theme["green"])
            if subtitle:
                add_textbox(slide, 0.68, 0.92, 11.6, 0.35, subtitle, 12, False, theme["muted"])
            line = slide.shapes.add_shape(1, Inches(0.65), Inches(1.28), Inches(12), Inches(0.03))
            line.fill.solid()
            line.fill.fore_color.rgb = theme["gold"]
            line.line.color.rgb = theme["gold"]

        def add_bullets(slide, left, top, width, height, lines, size=20):
            box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            frame = box.text_frame
            frame.clear()
            for idx, line in enumerate(lines):
                p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
                p.text = str(line)
                p.level = 0
                p.font.name = "Microsoft JhengHei"
                p.font.size = Pt(size)
                p.font.color.rgb = theme["ink"]

        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, theme["cream"])
        add_textbox(slide, 0.8, 1.35, 11.8, 0.8, data.get("title", "臺語教學簡報"), 34, True, theme["green"])
        add_textbox(
            slide,
            0.85,
            2.25,
            11.5,
            0.45,
            f"{data.get('grade', '')}｜{data.get('duration_minutes', '')} 分鐘",
            18,
            False,
            theme["gold"],
        )
        add_textbox(slide, 0.85, 5.9, 11.5, 0.35, "依據官方教材索引與 108 課綱輔助生成，教師審核後使用。", 13, False, theme["muted"])

        # 學習目標
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, RGBColor(255, 255, 255))
        add_title(slide, "學習目標與課綱連結", "用於課堂導入與教師備課確認")
        curriculum = data.get("curriculum", {}) or {}
        goal_lines = []
        goal_lines.extend(curriculum.get("learning_performance", [])[:3])
        goal_lines.extend(curriculum.get("learning_content", [])[:3])
        if not goal_lines:
            goal_lines = ["能理解本單元核心詞彙。", "能在情境中聽辨並使用臺語句型。", "能完成課堂互動測驗。"]
        add_bullets(slide, 0.9, 1.65, 11.5, 4.8, goal_lines[:6], 18)

        # 詞彙
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, RGBColor(255, 255, 255))
        add_title(slide, "核心詞彙", "漢字、臺羅與華語意義對照")
        vocab_lines = []
        for item in data.get("vocabulary", [])[:6]:
            vocab_lines.append(
                f"{item.get('hanji', '')}｜{item.get('tailo_diacritic') or item.get('tailo_numeric', '')}｜{item.get('zh_tw', '')}"
            )
        add_bullets(slide, 0.9, 1.65, 11.5, 4.8, vocab_lines or ["尚無詞彙資料"], 21)

        # 對話
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, RGBColor(255, 255, 255))
        add_title(slide, "情境會話", "可搭配聽說練習與角色扮演")
        dialogue_lines = []
        for item in data.get("dialogues", [])[:4]:
            dialogue_lines.append(f"{item.get('role', '')}：{item.get('hanji', '')}")
            if item.get("tailo_diacritic"):
                dialogue_lines.append(f"  {item.get('tailo_diacritic')}")
        add_bullets(slide, 0.9, 1.55, 11.5, 5.3, dialogue_lines or ["尚無對話資料"], 18)

        # 測驗
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, RGBColor(255, 255, 255))
        add_title(slide, "課堂測驗", "可轉作紙本考卷或互動網站題目")
        quiz_lines = []
        for idx, question in enumerate(data.get("questions", [])[:3], start=1):
            answer_index = question.get("answer_index", 0)
            options = question.get("options", []) or []
            answer = options[answer_index] if 0 <= answer_index < len(options) else ""
            quiz_lines.append(f"{idx}. {question.get('question', '')}")
            if answer:
                quiz_lines.append(f"   答案：{answer}")
        add_bullets(slide, 0.9, 1.55, 11.5, 5.2, quiz_lines or ["尚無測驗題目"], 17)

        official_slide_assets = data.get("official_slide_seeds", []) or []
        for asset in official_slide_assets[:2]:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_bg(slide, RGBColor(255, 255, 255))
            add_title(slide, "官方教材延伸素材", asset.get("heading") or asset.get("title", ""))
            source = asset.get("title", "")
            bullets = asset.get("bullets", []) or []
            lines = bullets[:5]
            if source:
                lines.append(f"來源：{source}")
            add_bullets(slide, 0.9, 1.55, 11.5, 5.2, lines or ["尚無官方延伸素材"], 16)

        # 官方教材來源
        official_materials = self._flatten_recommended_materials(data)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, theme["cream"])
        add_title(slide, "官方教材參考", "後續產出教材時可追溯來源")
        source_lines = []
        for item in official_materials[:5]:
            source_lines.append(self._format_official_material_line(item))
        add_bullets(slide, 0.9, 1.55, 11.5, 4.9, source_lines or ["本單元尚未命中官方教材索引，請教師補充來源。"], 17)

        pptx_path = os.path.join(output_dir, "teaching_slides.pptx")
        prs.save(pptx_path)
        print(f"  [+] 已產生 PowerPoint 簡報: {pptx_path}")
        
    def _generate_html(self, data: Dict[str, Any], output_dir: str):
        # 使用 Jinja2 樣板渲染離線互動網頁（樣板：templates/interactive_website.html.j2）
        from jinja2 import Environment, FileSystemLoader

        def js_str(s: Any) -> str:
            # 轉義字串以安全嵌入單引號 JS 字串字面值（避免漢字／翻譯內含單引號或反斜線破壞 onclick）
            return (str(s or "")
                    .replace(chr(92), chr(92)*2)
                    .replace("'", chr(92)+"'")
                    .replace(chr(10), " ")
                    .replace(chr(13), " "))

        template_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "templates")
        env = Environment(
            loader=FileSystemLoader(template_dir),
            autoescape=False,          # 內容含 JS 樣板字串與 HTML 片段，沿用原本不做 HTML escape 的行為
            trim_blocks=True,
            lstrip_blocks=True,
        )
        env.filters["js"] = js_str

        template = env.get_template("interactive_website.html.j2")
        curriculum = data.get("curriculum", {}) or {}
        html_output = template.render(
            title=data.get("title", ""),
            grade=data.get("grade", ""),
            duration_minutes=data.get("duration_minutes", ""),
            curriculum={
                "learning_performance": curriculum.get("learning_performance", []),
                "learning_content": curriculum.get("learning_content", []),
            },
            vocabulary=data.get("vocabulary", []),
            dialogues=data.get("dialogues", []),
            questions=data.get("questions", []),
            official_generation_assets=data.get("official_interactive_assets", []),
            official_materials=self._flatten_recommended_materials(data, limit=5),
        )

        html_path = os.path.join(output_dir, "interactive_website.html")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_output)
        print(f"  [+] 已產生離線網頁 HTML: {html_path}")

    def _generate_review_report(self, data: Dict[str, Any], output_dir: str):
        report_lines = [
            f"# 臺語教師教學審核報告：{data.get('title')}",
            f"產出日期: {datetime.now().strftime('%Y-%m-%d')}",
            "",
            "## 1. 課綱與教材檢核狀態",
            "本教材已串接 108 課綱與 seed 詞庫做字元比對，檢核結果如下：",
            ""
        ]
        
        report_lines.append("### 詞彙審查清單 (Vocabulary Review List)")
        for vocab in data.get("vocabulary", []):
            status = "✅ 已驗證" if vocab.get("review_status") == "teacher_verified" else "⚠️ 待教師確認"
            report_lines.append(f"- **{vocab.get('hanji')}** (`{vocab.get('tailo_diacritic')}`): {vocab.get('zh_tw')} -> **{status}**")
            
        report_lines.append("")
        report_lines.append("### 對話拼音聲調審查 (Dialogue Tone Review)")
        for dia in data.get("dialogues", []):
            report_lines.append(f"- **{dia.get('role')}**：{dia.get('hanji')} -> **待試聽音檔並校對**")

        # 自動內容檢核結果（華語用字、漢字/臺羅音節一致性）
        report_lines.append("")
        report_lines.append("## 2. 內容自動檢核 (Automated Content Check)")
        warnings = check_lesson_content(data)
        if warnings:
            report_lines.append(f"系統偵測到 {len(warnings)} 項待教師確認事項（僅為提醒，未自動更改）：")
            report_lines.append("")
            for w in warnings:
                report_lines.append(f"- ⚠️ {w}")
        else:
            report_lines.append("✅ 未發現明顯華語夾雜或漢字/臺羅音節不符問題。")

        report_lines.append("")
        report_lines.append("## 3. 官方教材素材建議")
        recommendations = data.get("official_material_recommendations", {}) or {}
        output_labels = {
            "exam": "考卷",
            "worksheet": "講義／學習單",
            "slides": "簡報",
            "video": "影片",
            "interactive": "互動網站",
            "quiz": "測驗",
        }
        if recommendations:
            for output, items in recommendations.items():
                report_lines.append("")
                report_lines.append(f"### {output_labels.get(output, output)}")
                if not items:
                    report_lines.append("- 尚未命中適合素材，建議教師補充來源。")
                    continue
                for item in items:
                    line = self._format_official_material_line(item)
                    url = item.get("page_url") or item.get("attachment_url") or ""
                    report_lines.append(f"- {line}  {url}")
        else:
            for item in self._flatten_recommended_materials(data):
                line = self._format_official_material_line(item)
                url = item.get("page_url") or item.get("attachment_url") or ""
                report_lines.append(f"- {line}  {url}")

        generation_assets = data.get("official_generation_assets", {}) or {}
        if generation_assets:
            report_lines.append("")
            report_lines.append("## 4. 官方教材可轉用生成素材")
            for output, assets in generation_assets.items():
                report_lines.append("")
                report_lines.append(f"### {output_labels.get(output, output)}")
                if not assets:
                    report_lines.append("- 尚未命中可直接轉用素材。")
                    continue
                for asset in assets:
                    asset_type = asset.get("asset_type", "")
                    title = asset.get("title", "")
                    if asset_type == "multiple_choice":
                        preview = asset.get("question", "")
                        note = "需教師確認正確答案" if asset.get("teacher_answer_required") else ""
                    elif asset_type == "vocabulary_card":
                        preview = f"{asset.get('hanji', '')}｜{asset.get('tailo', '')}｜{asset.get('meaning', '')}"
                        note = ""
                    elif asset_type == "reflection_4f":
                        preview = "；".join((asset.get("prompts", []) or [])[:2])
                        note = ""
                    else:
                        preview = "；".join((asset.get("bullets", []) or [])[:2])
                        note = ""
                    suffix = f"（{note}）" if note else ""
                    report_lines.append(f"- {asset_type}｜{title}：{preview}{suffix}")

        report_path = os.path.join(output_dir, "teacher_review_report.md")
        with open(report_path, "w", encoding="utf-8") as f:
            f.write("\n".join(report_lines))
        print(f"  [+] 已產生教師審核報告 Markdown: {report_path}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--config", default="config.json")
    parser.add_argument("--case", required=True)
    parser.add_argument("--output", default=None,
                        help="輸出目錄；未指定時讀 config 的 output.base_dir，再退回 \"output\"")
    parser.add_argument("--no-media", action="store_true",
                        help="跳過音訊與圖片生成，快速輸出考卷、講義、簡報與互動網頁")
    args = parser.parse_args()

    generator = MaterialGenerator(args.config)
    generator.generate_all(args.case, output_dir=args.output, skip_media=args.no_media)
