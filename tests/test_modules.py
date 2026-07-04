# 核心模組單元測試 (test_modules.py)
import os
import sys
import json
import pytest

# 加入 src 到搜尋路徑
sys.path.append(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "src"))

from tailo.validator import tailo_numeric_to_diacritic, convert_sentence_numeric_to_diacritic
from rag.retriever import TaigiRetriever

def test_tailo_conversion():
    # 測試聲調轉換
    assert tailo_numeric_to_diacritic("tsiah8") == "tsia̍h"
    assert tailo_numeric_to_diacritic("png7") == "pn̄g"
    assert tailo_numeric_to_diacritic("tshai3") == "tshài"
    assert tailo_numeric_to_diacritic("tshi7") == "tshī"
    assert tailo_numeric_to_diacritic("a2") == "á"
    assert tailo_numeric_to_diacritic("Tai5") == "Tâi"
    assert tailo_numeric_to_diacritic("uan5") == "uân"
    assert tailo_numeric_to_diacritic("thak8") == "tha̍k"
    assert tailo_numeric_to_diacritic("tsheh4") == "tsheh"
    assert tailo_numeric_to_diacritic("iu5") == "iû"

def test_sentence_conversion():
    # 測試整句轉換
    sentence_num = "tsiah8-png7 e7-poo khī-tshai3-tshi7-a2"
    sentence_dia = convert_sentence_numeric_to_diacritic(sentence_num)
    # 驗證拼音轉換 (部分詞語如 e7-poo, khī 未在單元測試字庫中但仍應套用轉換規則)
    assert "tsia̍h-pn̄g" in sentence_dia
    assert "tshài-tshī-á" in sentence_dia

def test_retriever():
    retriever = TaigiRetriever()
    
    # 測試詞彙檢索
    results = retriever.retrieve_vocabulary("食飯")
    assert len(results) > 0
    assert results[0]["hanji"] == "食飯"
    assert results[0]["tailo_numeric"] == "tsiah8-png7"
    
    # 測試課綱年級過濾 (國中)
    syllabus_jh = retriever.retrieve_syllabus_by_grade("國中七年級")
    assert len(syllabus_jh) > 0

    # 單字短詞不可被模糊替換成不相干長詞，例如「手」不應變成「手機」。
    assert all(item["hanji"] != "手機" for item in retriever.retrieve_vocabulary("手"))
    # 確保含有國中七至九年級課綱
    assert any("1-Ⅳ-1" in item["code"] for item in syllabus_jh)

def test_enrichment():
    retriever = TaigiRetriever()
    
    # 測試教材結構豐富化
    lesson_draft = {
        "title": "菜市場買物件",
        "grade": "國中七年級",
        "vocabulary": ["食飯", "偌濟錢"],
        "dialogues": []
    }
    
    enriched = retriever.enrich_lesson_json(lesson_draft)
    
    # 1. 確保詞彙成功被豐富為物件
    assert isinstance(enriched["vocabulary"][0], dict)
    assert enriched["vocabulary"][0]["hanji"] == "食飯"
    assert enriched["vocabulary"][0]["tailo_numeric"] == "tsiah8-png7"
    
    # 2. 確保課綱自動被帶出
    assert len(enriched["curriculum"]["learning_performance"]) > 0
    assert any("1-Ⅳ-1" in x for x in enriched["curriculum"]["learning_performance"])

@pytest.mark.network
def test_tts_fetch_moedict(tmp_path):
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    # 測試下載「菜市仔」
    output_file = tmp_path / "vocab_菜市仔.ogg"
    success = tts.fetch_vocab_audio("菜市仔", str(output_file))
    assert success is True
    assert output_file.exists()
    assert output_file.stat().st_size > 0

def test_tts_dummy_synthesize(tmp_path):
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    # 強制設定為 dummy
    tts.provider = "dummy"
    output_file = tmp_path / "sentence.wav"
    success = tts.synthesize_sentence("這是一句測試句子", str(output_file))
    assert success is True
    assert output_file.exists()
    assert output_file.stat().st_size > 0

def test_content_checker_flags_and_passes():
    """內容檢核：壞資料應被標記，修正後的 Mock 課文應零警告。"""
    from agent.content_checker import check_lesson_content
    bad = {
        "vocabulary": [{"hanji": "老闆"}],
        "dialogues": [{"hanji": "你欲去買一些物件嗎？", "tailo_numeric": "li2 beh4 khi3"}],
    }
    warns = check_lesson_content(bad)
    joined = " ".join(warns)
    assert any("老闆" in w for w in warns)      # 華語用字
    assert any("一些" in w for w in warns)
    assert "音節" in joined                       # 漢字/臺羅音節數不符

    # 修正後的 Mock（菜市分支）應無警告
    from agent.outline_generator import TaigiOutlineGenerator
    gen = TaigiOutlineGenerator()
    gen._get_available_models = lambda: []
    outline = gen.generate_outline("去菜市仔買物件", "國中七年級")
    assert check_lesson_content(outline) == []


def test_tailo_to_poj_conversion():
    """臺羅數字調 → 白話字調符式（供本地 mms-tts-nan 使用）。"""
    from tailo.poj import tailo_to_poj
    cases = {
        "tsiah8-png7": "chia̍h-pn̄g",       # 食飯：ts->ch、入聲調
        "tshai3-tshi7-a2": "chhài-chhī-á",  # 菜市仔：tsh->chh
        "be2-mih8-kiann7": "bé-mi̍h-kiāⁿ",  # 買物件：鼻化 -nn->ⁿ
        "gua7-tse7-tsinn5": "gōa-chē-chîⁿ", # 偌濟錢：oa 調符在 o、鼻化
        "thak8-tsheh4": "tha̍k-chheh",       # 讀冊
        "kue2-tsi2": "kóe-chí",             # 果子：ue->oe 且調符在 o
        "gua2": "góa",                      # 我：oa 調符在 o
        "ue7": "ōe",                        # 話：oe 調符在 o
        "kuai2": "koái",                    # oai 三母音仍標第二母音(a)
        "kong1-hng5": "kong-hn̂g",           # 公園：無聲母韻母變化
    }
    for num, expect in cases.items():
        assert tailo_to_poj(num) == expect, f"{num} -> {tailo_to_poj(num)} (期望 {expect})"


def test_mms_provider_falls_back_without_torch(tmp_path, monkeypatch):
    """未安裝 torch 時 mms provider 應安全降級，不得拋例外。"""
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    tts.provider = "mms"
    # 攔截降級目標，避免實際連網/ffmpeg
    monkeypatch.setattr(tts, "_synthesize_concatenative",
                        lambda text, out: tts._synthesize_dummy(text, out))
    out = tmp_path / "m.wav"
    # 有 tailo 但無 torch（測試環境未裝）→ 降級；無論如何都應產生有效檔
    ok = tts.synthesize_sentence("食飯", str(out), "tsiah8-png7")
    assert ok is True and out.exists() and out.stat().st_size > 0


def test_mms_provider_requires_tailo(tmp_path, monkeypatch):
    """mms provider 缺臺羅拼音時應改走接音合成路徑。"""
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    tts.provider = "mms"
    called = {"concat": False}
    def fake_concat(text, out):
        called["concat"] = True
        return tts._synthesize_dummy(text, out)
    monkeypatch.setattr(tts, "_synthesize_concatenative", fake_concat)
    out = tmp_path / "m.wav"
    tts.synthesize_sentence("食飯", str(out))  # 無 tailo_numeric
    assert called["concat"] is True


def test_tts_concat_segmentation(monkeypatch):
    """接音合成斷詞：最長詞優先，缺字略過（離線，mock 萌典查詢）。"""
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    available = {"菜市仔": "p1", "買": "p2", "物件": "p3", "錢": "p4"}
    monkeypatch.setattr(tts, "_get_word_audio", lambda w: available.get(w))

    # 最長詞優先：菜市仔 應整詞匹配，而非拆成 菜/市/仔
    assert [w for w, _ in tts._segment_for_audio("菜市仔")] == ["菜市仔"]
    # 詞組可拆為已收錄的詞
    assert [w for w, _ in tts._segment_for_audio("買物件")] == ["買", "物件"]
    # 含標點與缺字：x 無音檔應被略過
    assert [w for w, _ in tts._segment_for_audio("買x錢")] == ["買", "錢"]


def test_tts_concat_falls_back_to_dummy(tmp_path, monkeypatch):
    """接音合成在句中無任何可發音詞時，應降級為 dummy 並仍產生有效 wav。"""
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    tts.provider = "concat"
    monkeypatch.setattr(tts, "_get_word_audio", lambda w: None)  # 全部查無
    out = tmp_path / "s.wav"
    assert tts.synthesize_sentence("無收錄內容", str(out)) is True
    assert out.exists() and out.stat().st_size > 0

def test_voxcpm_feed_text_hanji_tailo(tmp_path):
    from tts.generator import TaigiTTS

    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({
        "tts": {
            "provider": "voxcpm",
            "voxcpm": {"input_mode": "hanji_tailo"},
        }
    }, ensure_ascii=False), encoding="utf-8")

    tts = TaigiTTS(str(config_path))
    assert tts._voxcpm_feed_text("菜蔬", "Tshài-se") == "菜蔬；臺羅：Tshài-se"


def test_stt_dummy(tmp_path):
    from stt.generator import TaigiSTT
    stt = TaigiSTT()
    stt.provider = "dummy"
    
    # 建立一個測試用的虛擬檔案
    audio_file = tmp_path / "test.webm"
    audio_file.write_bytes(b"dummy")
    
    recognized = stt.speech_to_text(str(audio_file), target_text="買物件")
    assert recognized == "買物件"

def test_api_stt_endpoint(tmp_path, monkeypatch):
    from fastapi.testclient import TestClient
    from server import app
    from stt.generator import TaigiSTT
    
    # Mock stt.speech_to_text 避免 real ASR 在 API 測試中干擾
    monkeypatch.setattr(TaigiSTT, "speech_to_text", lambda self, p, t="": t)
    
    client = TestClient(app)
    
    # 建立一個測試用的音訊檔案
    test_file = tmp_path / "test.webm"
    test_file.write_bytes(b"dummy audio content")
    
    with open(test_file, "rb") as f:
        response = client.post(
            "/api/stt",
            files={"file": ("test.webm", f, "audio/webm")},
            data={"target_text": "多謝"}
        )
        
    assert response.status_code == 200
    assert response.json() == {"text": "多謝"}

def test_outline_generator_mock():
    from agent.outline_generator import TaigiOutlineGenerator
    generator = TaigiOutlineGenerator()
    # 關閉 Ollama 模擬離線
    generator._get_available_models = lambda: []
    
    outline = generator.generate_outline("去菜市仔買物件", "國中七年級")
    assert outline["title"] == "情境會話：去菜市仔買物件"
    assert outline["grade"] == "國中七年級"
    assert "菜市仔" in outline["vocabulary"]
    assert len(outline["dialogues"]) > 0
    assert len(outline["questions"]) == 3

def test_outline_generator_choose_model():
    from agent.outline_generator import TaigiOutlineGenerator
    generator = TaigiOutlineGenerator()
    
    # 1. 測試預設模型存在時，選擇預設模型
    generator.configured_model = "SARC-Taigi-LLM-12b:latest"
    generator._get_available_models = lambda: ["SARC-Taigi-LLM-12b:latest", "gemma4:12b"]
    assert generator._choose_model() == "SARC-Taigi-LLM-12b:latest"
    
    # 2. 測試預設模型不存在時，自動選擇本地替代大模型 (例如 gemma4)
    generator._get_available_models = lambda: ["gemma4:12b", "qwen2.5-coder:1.5b"]
    assert generator._choose_model() == "gemma4:12b"

def test_outline_generator_passes_official_materials_to_ollama(monkeypatch):
    from agent.outline_generator import TaigiOutlineGenerator
    import agent.outline_generator as outline_module

    class FakeRetriever:
        def retrieve_vocabulary(self, prompt):
            return []

        def retrieve_syllabus_by_grade(self, grade):
            return []

        def retrieve_official_materials(self, prompt, limit=3):
            return [{
                "title": "官方有機菜蔬教材",
                "attachment_label": "資源頁",
                "learning_stage": "第五學習階段",
                "snippet": "官方教材片段：臺南佳里的有機菜蔬。"
            }]

    class FakeResponse:
        status_code = 200

        def json(self):
            return {
                "message": {
                    "content": json.dumps({
                        "title": "官方教材測試",
                        "grade": "國中七年級",
                        "duration_minutes": 45,
                        "vocabulary": ["菜蔬", "佳里"],
                        "dialogues": [{
                            "role": "老師",
                            "hanji": "咱來學菜蔬。",
                            "tailo_numeric": "lan2 lai5 oh8 tshai3-se1.",
                            "zh_tw": "我們來學蔬菜。"
                        }],
                        "questions": [{
                            "id": "q1",
                            "question": "這份教材主題是什麼？",
                            "options": ["菜蔬", "交通", "天氣", "動物"],
                            "answer_index": 0,
                            "explanation": "官方教材片段提到有機菜蔬。"
                        }]
                    }, ensure_ascii=False)
                }
            }

    captured = {}

    def fake_post(url, json, timeout):
        captured["system_prompt"] = json["messages"][0]["content"]
        return FakeResponse()

    generator = TaigiOutlineGenerator()
    generator.retriever = FakeRetriever()
    generator._choose_model = lambda: "fake-model"
    monkeypatch.setattr(outline_module.requests, "post", fake_post)

    outline = generator.generate_outline("有機菜蔬", "國中七年級")

    assert outline["title"] == "官方教材測試"
    assert "官方教材片段" in captured["system_prompt"]
    assert "臺南佳里的有機菜蔬" in captured["system_prompt"]

def test_natural_language_request_parser():
    from agent.natural_language_runner import parse_natural_language_request

    parsed = parse_natural_language_request("幫我做國中七年級有機菜蔬的考卷、簡報、互動網站和測驗")

    assert parsed["grade"] == "國中七年級"
    assert parsed["topic"] == "有機菜蔬"
    assert "exam" in parsed["outputs"]
    assert "slides" in parsed["outputs"]
    assert "interactive" in parsed["outputs"]
    assert "quiz" in parsed["outputs"]

def test_natural_language_runner_manifest_includes_recommendations(monkeypatch, tmp_path):
    import agent.natural_language_runner as runner

    class FakeRetriever:
        def recommend_official_materials(self, topic, outputs, limit_per_output=3):
            return {
                "video": [{
                    "title": "官方影片",
                    "resource_kind": "video",
                    "page_url": "https://mhi.moe.edu.tw/video",
                }]
            }

        def recommend_official_generation_assets(self, topic, outputs, limit_per_output=3):
            return {
                "video": [{
                    "asset_id": "asset-1",
                    "asset_type": "slide_seed",
                    "title": "官方生成素材",
                    "bullets": ["官方素材重點"],
                }]
            }

    class FakeOutlineGenerator:
        def __init__(self, config_path):
            self.retriever = FakeRetriever()

        def generate_outline(self, topic, grade, duration_minutes=45):
            return {
                "title": topic,
                "grade": grade,
                "duration_minutes": duration_minutes,
                "vocabulary": [],
                "dialogues": [],
                "questions": [],
            }

    class FakeMaterialGenerator:
        def __init__(self, config_path):
            pass

        def generate_all(self, outline_path, output_dir=None, skip_media=False):
            assert skip_media is True
            marker = tmp_path / "lesson_structure.json"
            marker.write_text("{}", encoding="utf-8")
            return str(tmp_path)

    monkeypatch.setattr(runner, "TaigiOutlineGenerator", FakeOutlineGenerator)
    monkeypatch.setattr(runner, "MaterialGenerator", FakeMaterialGenerator)

    manifest = runner.run_natural_language_request(
        "幫我做國中七年級有機菜蔬的影片",
        output_dir=str(tmp_path),
        include_video=False,
        skip_media=True,
        validate_output=False,
    )

    assert manifest["generation_options"]["skip_media"] is True
    assert manifest["official_material_recommendations"]["video"][0]["title"] == "官方影片"
    assert manifest["official_generation_assets"]["video"][0]["asset_type"] == "slide_seed"
    manifest_path = tmp_path / "generation_manifest.json"
    saved_manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert saved_manifest["generation_options"]["skip_media"] is True
    assert saved_manifest["official_material_recommendations"]["video"][0]["resource_kind"] == "video"
    assert saved_manifest["official_generation_assets"]["video"][0]["title"] == "官方生成素材"


def test_natural_language_runner_records_successful_video(monkeypatch, tmp_path):
    from pathlib import Path
    import agent.natural_language_runner as runner

    class FakeRetriever:
        def recommend_official_materials(self, topic, outputs, limit_per_output=3):
            return {"video": [{"title": "官方影片", "resource_kind": "video"}]}

        def recommend_official_generation_assets(self, topic, outputs, limit_per_output=3):
            return {"video": [{"asset_type": "slide_seed", "title": "官方分鏡"}]}

    class FakeOutlineGenerator:
        def __init__(self, config_path):
            self.retriever = FakeRetriever()

        def generate_outline(self, topic, grade, duration_minutes=45):
            return {
                "title": topic,
                "grade": grade,
                "duration_minutes": duration_minutes,
                "vocabulary": [],
                "dialogues": [],
                "questions": [],
            }

    class FakeMaterialGenerator:
        def __init__(self, config_path):
            pass

        def generate_all(self, outline_path, output_dir=None, skip_media=False):
            lesson_path = tmp_path / "lesson_structure.json"
            lesson_path.write_text("{}", encoding="utf-8")
            return str(tmp_path)

    class FakeVideoGenerator:
        def __init__(self, config_path):
            pass

        def generate_video(self, lesson_json_path, output_mp4_path):
            Path(output_mp4_path).write_bytes(b"fake mp4")
            return True

    monkeypatch.setattr(runner, "TaigiOutlineGenerator", FakeOutlineGenerator)
    monkeypatch.setattr(runner, "MaterialGenerator", FakeMaterialGenerator)
    monkeypatch.setattr(runner, "TaigiVideoGenerator", FakeVideoGenerator)

    manifest = runner.run_natural_language_request(
        "幫我做國中七年級有機菜蔬的影片",
        output_dir=str(tmp_path),
        include_video=None,
        skip_media=True,
        validate_output=False,
    )

    assert manifest["generation_options"]["include_video"] is True
    assert manifest["video_generation"]["attempted"] is True
    assert manifest["video_generation"]["success"] is True
    assert manifest["outputs"]["video"].endswith("lesson_video.mp4")
    assert (tmp_path / "lesson_video.mp4").read_bytes() == b"fake mp4"


def test_outline_generator_mock_default():
    from agent.outline_generator import TaigiOutlineGenerator
    generator = TaigiOutlineGenerator()
    generator._get_available_models = lambda: []
    
    outline = generator.generate_outline("在茶莊品茶", "國中七年級")
    assert outline["title"] == "情境會話：在茶莊品茶"
    assert outline["grade"] == "國中七年級"
    assert len(outline["vocabulary"]) == 3
    assert "食飯" in outline["vocabulary"]
    assert len(outline["dialogues"]) == 3
    assert len(outline["questions"]) == 3  # 驗證已補上第 3 題
    assert outline["questions"][2]["id"] == "q3"

def test_outline_generator_mock_body_topic():
    from agent.outline_generator import TaigiOutlineGenerator
    generator = TaigiOutlineGenerator()
    generator._get_available_models = lambda: []

    outline = generator.generate_outline("身體五官", "國中七年級")
    vocab_words = [item["hanji"] if isinstance(item, dict) else item for item in outline["vocabulary"]]

    assert "身軀" in vocab_words
    assert "目睭" in vocab_words
    assert any("五官" in dialogue["hanji"] for dialogue in outline["dialogues"])


def test_outline_generator_mock_vegetable_topic():
    from agent.outline_generator import TaigiOutlineGenerator
    generator = TaigiOutlineGenerator()
    generator._get_available_models = lambda: []

    outline = generator.generate_outline("有機菜蔬", "國中七年級")
    vocab_words = [item["hanji"] if isinstance(item, dict) else item for item in outline["vocabulary"]]

    assert "菜蔬" in vocab_words
    assert "有機" in vocab_words
    assert any("菜蔬" in question["question"] or "菜蔬" in question["explanation"] for question in outline["questions"])

def test_outline_generator_compile_flag(monkeypatch, tmp_path):
    from agent.outline_generator import TaigiOutlineGenerator
    generator = TaigiOutlineGenerator()
    generator._get_available_models = lambda: []
    
    outline = generator.generate_outline("去菜市仔買物件", "國中七年級")
    assert outline["vocabulary"][0] == "菜市仔"
    assert len(outline["dialogues"]) == 4

def test_image_generator_config_url():
    from generators.image_generator import FreeImageGenerator
    import json, tempfile, os
    
    config = {"ollama": {"url": "http://localhost:9999", "model": "test-model"}}
    tmp = tempfile.NamedTemporaryFile(mode="w", suffix=".json", delete=False)
    json.dump(config, tmp)
    tmp.close()
    
    try:
        gen = FreeImageGenerator(tmp.name)
        assert gen.ollama_url == "http://localhost:9999"
    finally:
        os.unlink(tmp.name)

def test_config_loaders_accept_utf8_bom(tmp_path):
    from agent.outline_generator import TaigiOutlineGenerator
    from generators.material_generator import MaterialGenerator
    from generators.image_generator import FreeImageGenerator
    from tts.generator import TaigiTTS

    config_path = tmp_path / "config.json"
    config_bytes = b"\xef\xbb\xbf" + json.dumps({
        "ollama": {"url": "http://localhost:9999", "model": "fake"},
        "tts": {"provider": "dummy"},
        "piauim": {"provider": "off"},
    }).encode("utf-8")
    config_path.write_bytes(config_bytes)

    assert TaigiOutlineGenerator(str(config_path)).config["ollama"]["model"] == "fake"
    assert MaterialGenerator(str(config_path)).config["piauim"]["provider"] == "off"
    assert FreeImageGenerator(str(config_path)).default_model == "fake"
    assert TaigiTTS(str(config_path)).provider == "dummy"

def test_sentence_conversion_mixed():
    from tailo.validator import convert_sentence_numeric_to_diacritic
    result = convert_sentence_numeric_to_diacritic("tsiah8-png7 to1-sia7")
    assert "tsia̍h-pn̄g" in result
    assert "to-siā" in result

def test_tailo_iou_ui_rules():
    from tailo.validator import tailo_numeric_to_diacritic
    assert tailo_numeric_to_diacritic("iu5") == "iû"
    assert tailo_numeric_to_diacritic("ui7") == "uī"
    assert tailo_numeric_to_diacritic("gue5") == "guê"
    assert tailo_numeric_to_diacritic("ng5") == "n̂g"
    assert tailo_numeric_to_diacritic("m7") == "m̄"

def test_retriever_no_results():
    from rag.retriever import TaigiRetriever
    retriever = TaigiRetriever()
    results = retriever.retrieve_vocabulary("不存在的詞彙xyz")
    assert len(results) == 0

def test_retriever_official_materials_index():
    from rag.retriever import TaigiRetriever
    retriever = TaigiRetriever()
    results = retriever.retrieve_official_materials("菜蔬", limit=3)
    assert len(results) > 0
    assert any("臺南佳里" in item["title"] for item in results)
    assert all(item["snippet"] for item in results)

    page_only_results = retriever.retrieve_official_materials("字音字形", limit=3)
    assert len(page_only_results) > 0
    assert any("字音字形" in item["title"] for item in page_only_results)

    interactive_results = retriever.retrieve_official_materials("wordwall", limit=3)
    assert any(item["resource_kind"] == "interactive" for item in interactive_results)
    assert any(
        link.get("domain") == "wordwall.net"
        for item in interactive_results
        for link in item.get("related_links", [])
    )

    video_results = retriever.retrieve_official_materials("有機菜蔬", limit=5)
    assert any(item["resource_kind"] == "video" for item in video_results)
    assert any(
        link.get("domain") in {"www.youtube.com", "youtu.be"}
        for item in video_results
        for link in item.get("related_links", [])
    )

    natural_query_results = retriever.retrieve_official_materials("身體五官互動遊戲", limit=3)
    assert natural_query_results[0]["resource_kind"] == "interactive"
    assert "身體五官" in natural_query_results[0]["title"]

    multi_term_video_results = retriever.retrieve_official_materials("台南佳里有機菜蔬", limit=3)
    assert multi_term_video_results[0]["resource_kind"] == "video"
    assert "有機菜蔬" in multi_term_video_results[0]["title"]

    recommendations = retriever.recommend_official_materials(
        "身體五官互動遊戲",
        outputs=["interactive", "quiz"],
        limit_per_output=2,
    )
    assert recommendations["interactive"][0]["resource_kind"] == "interactive"
    assert "身體五官" in recommendations["interactive"][0]["title"]

    video_recommendations = retriever.recommend_official_materials(
        "台南佳里有機菜蔬",
        outputs=["video"],
        limit_per_output=2,
    )
    assert video_recommendations["video"][0]["resource_kind"] == "video"
    assert "有機菜蔬" in video_recommendations["video"][0]["title"]

def test_retriever_official_material_snippets_index():
    from rag.retriever import TaigiRetriever
    retriever = TaigiRetriever()

    results = retriever.retrieve_official_material_snippets("有機菜蔬", limit=5)

    assert len(results) > 0
    assert all(item["resource_kind"] == "text_snippet" for item in results)
    assert any("有機菜蔬" in item["title"] or "有機菜蔬" in item["snippet"] for item in results)
    assert all(item["snippet"] for item in results)

def test_retriever_official_material_bank_index():
    from rag.retriever import TaigiRetriever
    retriever = TaigiRetriever()

    results = retriever.retrieve_official_material_bank(
        "有機菜蔬",
        material_kinds=["worksheet", "assessment"],
        limit=5,
    )

    assert len(results) > 0
    assert all(item["resource_kind"] == "material_bank" for item in results)
    assert any("worksheet" in item["material_kinds"] for item in results)
    assert any("有機菜蔬" in item["title"] or "有機菜蔬" in item["snippet"] for item in results)
    assert any(item["curriculum_codes"] or "選擇題" in item["snippet"] for item in results)

def test_retriever_official_generation_assets_index():
    from rag.retriever import TaigiRetriever
    retriever = TaigiRetriever()

    results = retriever.retrieve_official_generation_assets(
        "有機菜蔬",
        output_tags=["exam"],
        asset_types=["multiple_choice"],
        limit=3,
    )

    assert len(results) >= 3
    assert all(item["resource_kind"] == "generation_asset" for item in results)
    assert all(item["asset_type"] == "multiple_choice" for item in results)
    assert any("有機" in item["question"] or "菜蔬" in " ".join(item["options"]) for item in results)

    recommendations = retriever.recommend_official_generation_assets(
        "有機菜蔬",
        outputs=["exam", "slides", "interactive", "quiz"],
        limit_per_output=2,
    )
    assert recommendations["exam"][0]["asset_type"] == "multiple_choice"
    assert any(item["asset_type"] == "slide_seed" for item in recommendations["slides"])

def test_official_catalog_analysis_outputs():
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "analyze_official_catalog.py"
    spec = importlib.util.spec_from_file_location("analyze_official_catalog", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    catalog_path = project_root / "data" / "official_materials" / "catalog.json"
    catalog = json.loads(catalog_path.read_text(encoding="utf-8-sig"))
    analysis = module.analyze(catalog)

    assert analysis["catalog_count"] == len(catalog)
    assert analysis["local_pdf_count"] >= 50
    assert any(row["name"] == "video" for row in analysis["by_resource_kind"])
    assert any(row["name"] == "食物與市場" for row in analysis["topic_coverage"])

def test_official_repository_audit_outputs():
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "audit_official_repository.py"
    spec = importlib.util.spec_from_file_location("audit_official_repository", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    data_root = project_root / "data" / "official_materials"
    catalog = json.loads((data_root / "catalog.json").read_text(encoding="utf-8-sig"))
    sources = json.loads((data_root / "sources.json").read_text(encoding="utf-8-sig"))
    pdf_text_index = json.loads((data_root / "analysis" / "pdf_text_index.json").read_text(encoding="utf-8-sig"))
    errors = module.load_jsonl(data_root / "analysis" / "mhi_collect_errors.jsonl")
    report = module.audit_repository(catalog, sources, pdf_text_index, errors, project_root)

    assert report["catalog_count"] == len(catalog)
    assert report["local_pdf_count"] >= 50
    assert report["pdf_text_missing_count"] == 0
    assert report["issue_counts"]["missing_local_files"] == 0
    assert report["issue_counts"]["zero_byte_local_files"] == 0
    assert report["issue_counts"]["source_count_mismatches"] == 0
    assert report["issue_counts"]["downloadable_without_local"] == 0
    assert report["issue_counts"]["confirmed_unavailable"] == 2
    assert any(row["source_id"] == "moe-mhi-108-minnan" for row in report["source_checks"])

def test_official_snippet_index_builder_outputs():
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "build_official_snippet_index.py"
    spec = importlib.util.spec_from_file_location("build_official_snippet_index", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    pdf_text_index_path = project_root / "data" / "official_materials" / "analysis" / "pdf_text_index.json"
    pdf_text_index = json.loads(pdf_text_index_path.read_text(encoding="utf-8-sig"))
    index = module.build_snippet_index(pdf_text_index, project_root)

    assert index["source_pdf_count"] == len(pdf_text_index)
    assert index["snippet_count"] >= 500
    assert index["skipped_count"] == 0
    assert any(row["name"] == "食物與市場" for row in index["by_topic"])
    assert any("有機菜蔬" in item["title"] or "有機菜蔬" in item["snippet"] for item in index["snippets"])

def test_official_material_bank_builder_outputs():
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "build_official_material_bank.py"
    spec = importlib.util.spec_from_file_location("build_official_material_bank", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    snippet_index_path = project_root / "data" / "official_materials" / "analysis" / "official_material_snippets.json"
    snippet_index = json.loads(snippet_index_path.read_text(encoding="utf-8-sig"))
    bank = module.build_material_bank(snippet_index)

    assert bank["source_snippet_count"] == snippet_index["snippet_count"]
    assert bank["bank_item_count"] >= 500
    assert bank["curriculum_item_count"] > 0
    assert any(row["name"] == "worksheet" for row in bank["by_material_kind"])
    assert any(row["name"] == "assessment" for row in bank["by_material_kind"])
    assert any("有機菜蔬" in item["title"] or "有機菜蔬" in item["snippet"] for item in bank["items"])

def test_official_generation_packs_builder_outputs():
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "build_official_generation_packs.py"
    spec = importlib.util.spec_from_file_location("build_official_generation_packs", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    bank_path = project_root / "data" / "official_materials" / "analysis" / "official_material_bank.json"
    bank = json.loads(bank_path.read_text(encoding="utf-8-sig"))
    packs = module.build_generation_packs(bank)

    assert packs["source_bank_item_count"] == bank["bank_item_count"]
    assert packs["asset_count"] >= 100
    assert any(row["name"] == "multiple_choice" for row in packs["by_asset_type"])
    assert any(row["name"] == "slide_seed" for row in packs["by_asset_type"])
    assert any(asset["asset_type"] == "multiple_choice" and "有機" in asset["question"] for asset in packs["assets"])

def test_free_image_generator(tmp_path, monkeypatch):
    from generators.image_generator import FreeImageGenerator
    import requests
    
    gen = FreeImageGenerator()
    # Mock 翻譯為 "apple"
    gen.translate_to_english_prompt = lambda x: "apple"
    
    # 建立 Mock 的 Response 物件
    class MockResponse:
        def __init__(self, json_data, status_code):
            self.json_data = json_data
            self.status_code = status_code
            self.content = b"fake image bytes"
            self.text = json.dumps(json_data)
            
        def json(self):
            return self.json_data
            
    # Mock post (提交任務)
    def mock_post(*args, **kwargs):
        return MockResponse({"id": "mock-task-123"}, 202)
        
    # Mock get (狀態查詢與下載)
    def mock_get(url, *args, **kwargs):
        if "check" in url:
            return MockResponse({"done": True}, 200)
        elif "status" in url:
            return MockResponse({
                "generations": [{"img": "https://example.com/mock.jpg"}]
            }, 200)
        else:
            # 下載圖片
            return MockResponse({}, 200)
            
    monkeypatch.setattr(requests, "post", mock_post)
    monkeypatch.setattr(requests, "get", mock_get)
    
    out_path = tmp_path / "vocab_apple.jpg"
    success = gen.generate_image("蘋果", str(out_path), width=128, height=128)
    assert success is True
    assert out_path.exists()
    assert out_path.read_bytes() == b"fake image bytes"

def test_free_image_generator_fallback_image(tmp_path, monkeypatch):
    from generators.image_generator import FreeImageGenerator
    from PIL import Image
    import requests

    gen = FreeImageGenerator()
    gen.translate_to_english_prompt = lambda x: "vegetables"

    def mock_post(*args, **kwargs):
        raise requests.RequestException("network down")

    monkeypatch.setattr(requests, "post", mock_post)
    monkeypatch.setattr(gen, "_http_post_with_fallback", lambda *args, **kwargs: (500, ""))

    out_path = tmp_path / "vocab_菜蔬.jpg"
    success = gen.generate_image("菜蔬", str(out_path), width=128, height=128)

    assert success is True
    assert out_path.exists()
    with Image.open(out_path) as img:
        assert img.size == (128, 128)

def test_video_generator_basic():
    from generators.video_generator import TaigiVideoGenerator
    gen = TaigiVideoGenerator()
    # 測試時長獲取預設值
    assert gen.get_audio_duration("") == 3.0
    if os.name == "nt":
        assert gen.npm_cmd.lower().endswith((".cmd", "npm"))
        assert not gen.npm_cmd.lower().endswith(".ps1")
        assert gen.npx_cmd.lower().endswith((".cmd", "npx"))
        assert not gen.npx_cmd.lower().endswith(".ps1")
    assert gen._has_command(gen.node_cmd)

@pytest.mark.network
def test_stt_local_whisper_integration(tmp_path):
    from stt.generator import TaigiSTT
    stt = TaigiSTT()
    # 強制設定為 whisper 模式以執行實體推理
    stt.provider = "whisper"
    stt.stt_config["local_model_size"] = "tiny"
    
    # 藉由 TaigiTTS 下載一個合法的台語語音檔 (例如「食飯」)
    from tts.generator import TaigiTTS
    tts = TaigiTTS()
    audio_file = tmp_path / "vocab_食飯.ogg"
    fetch_success = tts.fetch_vocab_audio("食飯", str(audio_file))
    
    if fetch_success:
        recognized = stt.speech_to_text(str(audio_file))
        assert len(recognized) > 0

def test_outline_validator_rejects_bad_romanization():
    from agent.outline_generator import TaigiOutlineGenerator
    gen = TaigiOutlineGenerator()
    
    bad_data = {
        "title": "測試",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "vocabulary": ["食飯", "買物件"],
        "dialogues": [
            {
                "role": "阿偉",
                "hanji": "你好",
                "tailo_numeric": "qiao2-lun4-ni3-san4",
                "zh_tw": "你好"
            }
        ],
        "questions": [
            {
                "id": "q1",
                "question": "測試",
                "options": ["A", "B", "C", "D"],
                "answer_index": 0,
                "explanation": "測試"
            }
        ]
    }
    assert gen._validate_outline(bad_data) is False

def test_outline_validator_rejects_grade_simplified():
    from agent.outline_generator import TaigiOutlineGenerator
    gen = TaigiOutlineGenerator()
    
    data = {
        "title": "測試",
        "grade": "七年级",
        "duration_minutes": 45,
        "vocabulary": ["食飯", "買物件"],
        "dialogues": [
            {
                "role": "阿偉",
                "hanji": "咱來去食飯",
                "tailo_numeric": "lan2 lai5-khi3 tsiah8-png7",
                "zh_tw": "我們來去吃飯"
            }
        ],
        "questions": [
            {
                "id": "q1",
                "question": "測試",
                "options": ["A", "B", "C", "D"],
                "answer_index": 0,
                "explanation": "測試"
            }
        ]
    }
    assert gen._validate_outline(data) is False

def test_outline_validator_rejects_duplicate_options():
    from agent.outline_generator import TaigiOutlineGenerator
    gen = TaigiOutlineGenerator()
    
    data = {
        "title": "測試",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "vocabulary": ["食飯", "買物件"],
        "dialogues": [
            {
                "role": "阿偉",
                "hanji": "咱來去食飯",
                "tailo_numeric": "lan2 lai5-khi3 tsiah8-png7",
                "zh_tw": "我們來去吃飯"
            }
        ],
        "questions": [
            {
                "id": "q1",
                "question": "測試",
                "options": ["相同", "相同", "相同", "相同"],
                "answer_index": 0,
                "explanation": "測試"
            }
        ]
    }
    assert gen._validate_outline(data) is False

def test_outline_validator_accepts_good_data():
    from agent.outline_generator import TaigiOutlineGenerator
    gen = TaigiOutlineGenerator()
    
    good_data = {
        "title": "情境會話：去菜市仔買物件",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "vocabulary": ["菜市仔", "買物件", "偌濟錢"],
        "dialogues": [
            {
                "role": "阿偉",
                "hanji": "阿媽，咱今仔日欲去菜市仔買物件無？",
                "tailo_numeric": "a1-ma2 lan2 kin1-a2-jit8 beh4 khi3 tshai3-tshi7-a2 be2-mih8-kiann7 bo5",
                "zh_tw": "阿嬤，我們今天要去菜市場買東西嗎？"
            }
        ],
        "questions": [
            {
                "id": "q1",
                "question": "測試題目",
                "options": ["選項A", "選項B", "選項C", "選項D"],
                "answer_index": 0,
                "explanation": "解析說明"
            }
        ]
    }
    assert gen._validate_outline(good_data) is True

def test_html_pronunciation_eval_not_broken(tmp_path):
    """回歸測試：確保產生的互動網頁 JS 樣板字串沒有多餘的 $（曾導致發音評估功能整個失效）。"""
    from generators.material_generator import MaterialGenerator
    gen = MaterialGenerator()

    data = {
        "title": "測試單元",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "curriculum": {"learning_performance": [], "learning_content": []},
        "vocabulary": [
            {"hanji": "食飯", "tailo_diacritic": "tsia̍h-pn̄g", "zh_tw": "吃飯",
             "audio_file": "audio/vocab_食飯.wav", "image_file": ""}
        ],
        "dialogues": [
            {"role": "阿偉", "hanji": "咱來去食飯", "tailo_diacritic": "lán lâi-khì tsia̍h-pn̄g",
             "zh_tw": "我們來去吃飯", "audio_file": "audio/dialogue_0.wav"}
        ],
        "questions": [
            {"id": "q1", "question": "測試", "options": ["A", "B", "C", "D"],
             "answer_index": 0, "explanation": "解析"}
        ]
    }

    gen._generate_html(data, str(tmp_path))
    html = (tmp_path / "interactive_website.html").read_text(encoding="utf-8")

    # 1. 絕不能再出現多餘的 $${ 樣式（這是先前的 bug）
    assert "$${" not in html
    # 2. 正確的 JS 樣板插值應原樣輸出
    assert "evaluateSpeech(event, '${type}', ${idx}, '${targetText}')" in html
    assert "playRecord(event, '${audioUrl}')" in html
    assert "${scoreClass}" in html


def test_html_handles_empty_role(tmp_path):
    """role 為空字串時不應拋 IndexError。"""
    from generators.material_generator import MaterialGenerator
    gen = MaterialGenerator()
    data = {
        "title": "測試", "grade": "國中七年級", "duration_minutes": 45,
        "curriculum": {"learning_performance": [], "learning_content": []},
        "vocabulary": [],
        "dialogues": [{"role": "", "hanji": "test", "tailo_diacritic": "test",
                       "zh_tw": "測試", "audio_file": ""}],
        "questions": []
    }
    gen._generate_html(data, str(tmp_path))  # 不應拋例外
    assert (tmp_path / "interactive_website.html").exists()


def test_html_includes_official_material_recommendations(tmp_path):
    from generators.material_generator import MaterialGenerator
    gen = MaterialGenerator()
    data = {
        "title": "測試", "grade": "國中七年級", "duration_minutes": 45,
        "curriculum": {"learning_performance": [], "learning_content": []},
        "vocabulary": [],
        "dialogues": [],
        "questions": [],
        "official_material_recommendations": {
            "interactive": [{
                "title": "wordwall (Game) 身體五官對對碰",
                "learning_stage": "第一學習階段",
                "resource_kind": "interactive",
                "provider": "wordwall",
                "page_url": "https://mhi.moe.edu.tw/example",
            }]
        },
    }
    gen._generate_html(data, str(tmp_path))
    html = (tmp_path / "interactive_website.html").read_text(encoding="utf-8")
    assert "官方素材建議" in html
    assert "wordwall (Game) 身體五官對對碰" in html


def test_html_includes_official_generation_assets(tmp_path):
    from generators.material_generator import MaterialGenerator
    gen = MaterialGenerator()
    data = {
        "title": "測試", "grade": "國中七年級", "duration_minutes": 45,
        "curriculum": {"learning_performance": [], "learning_content": []},
        "vocabulary": [],
        "dialogues": [],
        "questions": [],
        "official_interactive_assets": [{
            "asset_type": "multiple_choice",
            "question": "官方教材提到的菜蔬情境是啥？",
            "options": ["菜市仔", "車站", "海邊", "郵局"],
            "teacher_answer_required": True,
            "learning_stage": "第四學習階段",
            "title": "官方菜蔬教材",
            "source_url": "https://mhi.moe.edu.tw/example",
        }],
    }
    gen._generate_html(data, str(tmp_path))
    html = (tmp_path / "interactive_website.html").read_text(encoding="utf-8")
    assert "官方延伸互動素材" in html
    assert "官方教材提到的菜蔬情境是啥？" in html
    assert "需教師確認正確答案後使用" in html


def test_material_generator_pptx_output(tmp_path):
    from generators.material_generator import MaterialGenerator
    from pptx import Presentation

    gen = MaterialGenerator()
    data = {
        "title": "測試單元",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "curriculum": {
            "learning_performance": ["能聽懂核心詞彙"],
            "learning_content": ["情境會話"],
        },
        "vocabulary": [{
            "hanji": "菜蔬",
            "tailo_diacritic": "tshài-se",
            "zh_tw": "蔬菜",
        }],
        "dialogues": [{
            "role": "老師",
            "hanji": "咱來學菜蔬。",
            "tailo_diacritic": "lán lâi o̍h tshài-se",
        }],
        "questions": [{
            "question": "主題是啥？",
            "options": ["菜蔬", "天氣", "交通", "動物"],
            "answer_index": 0,
            "explanation": "本單元學菜蔬。",
        }],
        "official_materials": [{
            "title": "官方教材",
            "learning_stage": "第三學習階段",
            "resource_kind": "video",
        }],
    }

    gen._generate_pptx(data, str(tmp_path))
    pptx_path = tmp_path / "teaching_slides.pptx"
    assert pptx_path.exists()
    prs = Presentation(str(pptx_path))
    assert len(prs.slides) == 6


def test_material_generator_pptx_includes_official_slide_assets(tmp_path):
    from generators.material_generator import MaterialGenerator
    from pptx import Presentation

    gen = MaterialGenerator()
    data = {
        "title": "測試單元",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "curriculum": {"learning_performance": [], "learning_content": []},
        "vocabulary": [],
        "dialogues": [],
        "questions": [],
        "official_slide_seeds": [{
            "title": "官方菜蔬教材",
            "heading": "菜蔬情境",
            "bullets": ["觀察菜市仔用語", "比較臺語與華語說法"],
        }],
    }

    gen._generate_pptx(data, str(tmp_path))
    prs = Presentation(str(tmp_path / "teaching_slides.pptx"))
    text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert len(prs.slides) == 7
    assert "官方教材延伸素材" in text
    assert "觀察菜市仔用語" in text


def test_teacher_review_report_groups_official_recommendations(tmp_path):
    from generators.material_generator import MaterialGenerator
    gen = MaterialGenerator()
    data = {
        "title": "測試單元",
        "vocabulary": [],
        "dialogues": [],
        "official_material_recommendations": {
            "exam": [{
                "title": "官方學習單",
                "learning_stage": "第五學習階段",
                "resource_kind": "learning_resource",
                "page_url": "https://mhi.moe.edu.tw/exam",
            }],
            "video": [{
                "title": "官方影片",
                "learning_stage": "第五學習階段",
                "resource_kind": "video",
                "page_url": "https://mhi.moe.edu.tw/video",
            }],
        },
    }
    gen._generate_review_report(data, str(tmp_path))
    report = (tmp_path / "teacher_review_report.md").read_text(encoding="utf-8")
    assert "官方教材素材建議" in report
    assert "### 考卷" in report
    assert "官方學習單" in report
    assert "### 影片" in report
    assert "官方影片" in report


def test_material_generator_exam_output(tmp_path):
    from generators.material_generator import MaterialGenerator
    from docx import Document

    gen = MaterialGenerator()
    data = {
        "title": "測試單元",
        "grade": "國中七年級",
        "vocabulary": [{
            "hanji": "菜蔬",
            "tailo_diacritic": "tshài-se",
            "zh_tw": "蔬菜",
        }],
        "questions": [{
            "question": "主題是啥？",
            "options": ["菜蔬", "天氣", "交通", "動物"],
            "answer_index": 0,
            "explanation": "本單元學菜蔬。",
        }],
        "official_generated_questions": [{
            "asset_type": "multiple_choice",
            "question": "官方教材延伸題題幹？",
            "options": ["選項一", "選項二", "選項三", "選項四"],
            "title": "官方學習單",
            "teacher_answer_required": True,
        }],
    }

    gen._generate_exam_docx(data, str(tmp_path))

    exam_path = tmp_path / "exam_paper.docx"
    answer_path = tmp_path / "exam_answer_key.docx"
    assert exam_path.exists()
    assert answer_path.exists()
    assert "臺語單元考卷" in Document(str(exam_path)).paragraphs[0].text
    exam_text = "\n".join(p.text for p in Document(str(exam_path)).paragraphs)
    answer_text = "\n".join(p.text for p in Document(str(answer_path)).paragraphs)
    assert "官方教材延伸題" in exam_text
    assert "官方教材延伸題題幹？" in exam_text
    assert "菜蔬" in answer_text
    assert "需教師依官方教材確認" in answer_text


def test_material_generator_quiz_bank_output(tmp_path):
    from generators.material_generator import MaterialGenerator

    gen = MaterialGenerator()
    data = {
        "title": "測試單元",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "questions": [{
            "id": "q1",
            "question": "主題是啥？",
            "options": ["菜蔬", "天氣", "交通", "動物"],
            "answer_index": 0,
            "explanation": "本單元學菜蔬。",
        }],
        "official_generated_questions": [{
            "asset_id": "official-1",
            "asset_type": "multiple_choice",
            "question": "官方教材延伸題題幹？",
            "options": ["選項一", "選項二", "選項三", "選項四"],
            "title": "官方學習單",
            "source_url": "https://mhi.moe.edu.tw/example",
        }],
    }

    gen._generate_quiz_bank(data, str(tmp_path))

    quiz_bank = json.loads((tmp_path / "quiz_bank.json").read_text(encoding="utf-8"))
    teacher_key = (tmp_path / "quiz_teacher_key.md").read_text(encoding="utf-8")
    assert quiz_bank["question_count"] == 1
    assert quiz_bank["questions"][0]["answer"] == "菜蔬"
    assert quiz_bank["questions"][0]["graded"] is True
    assert quiz_bank["official_extension_count"] == 1
    assert quiz_bank["official_extension_questions"][0]["graded"] is False
    assert quiz_bank["official_extension_questions"][0]["teacher_answer_required"] is True
    assert "需教師依官方教材確認" in teacher_key


def test_material_generator_no_media_mode_skips_tts_and_images(tmp_path, monkeypatch):
    from generators.material_generator import MaterialGenerator

    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({
        "piauim": {"provider": "off"},
        "tts": {"provider": "dummy"},
        "output": {"base_dir": str(tmp_path / "out")},
    }, ensure_ascii=False), encoding="utf-8")

    case_path = tmp_path / "case.json"
    case_path.write_text(json.dumps({
        "title": "有機菜蔬",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "vocabulary": [{
            "hanji": "菜蔬",
            "tailo_numeric": "tshai3-se1",
            "zh_tw": "蔬菜",
        }],
        "dialogues": [{
            "role": "老師",
            "hanji": "咱來學菜蔬。",
            "tailo_numeric": "lan2 lai5 oh8 tshai3-se1.",
            "zh_tw": "我們來學蔬菜。",
        }],
        "questions": [{
            "id": "q1",
            "question": "主題是啥？",
            "options": ["菜蔬", "天氣", "交通", "動物"],
            "answer_index": 0,
            "explanation": "本單元學菜蔬。",
        }],
    }, ensure_ascii=False), encoding="utf-8")

    generator = MaterialGenerator(str(config_path))
    monkeypatch.setattr(generator.tts, "fetch_vocab_audio", lambda *args, **kwargs: (_ for _ in ()).throw(AssertionError("不應下載音檔")))
    monkeypatch.setattr(generator.tts, "synthesize_sentence", lambda *args, **kwargs: (_ for _ in ()).throw(AssertionError("不應合成語音")))
    monkeypatch.setattr(generator.tts, "synthesize_voxcpm_batch", lambda *args, **kwargs: (_ for _ in ()).throw(AssertionError("不應批次合成語音")))

    out_dir = tmp_path / "out"
    generator.generate_all(str(case_path), output_dir=str(out_dir), skip_media=True)

    lesson = json.loads((out_dir / "lesson_structure.json").read_text(encoding="utf-8"))
    assert lesson["generation_options"]["skip_media"] is True
    assert lesson["vocabulary"][0]["audio_file"] == ""
    assert lesson["vocabulary"][0]["image_file"] == ""
    assert lesson["dialogues"][0]["audio_file"] == ""
    assert (out_dir / "exam_paper.docx").exists()
    assert (out_dir / "quiz_bank.json").exists()
    assert (out_dir / "quiz_teacher_key.md").exists()
    assert (out_dir / "teaching_slides.pptx").exists()
    assert (out_dir / "interactive_website.html").exists()


def test_generation_output_validator_accepts_no_media_output(tmp_path, monkeypatch):
    import importlib.util
    from generators.material_generator import MaterialGenerator

    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({
        "piauim": {"provider": "off"},
        "tts": {"provider": "dummy"},
        "output": {"base_dir": str(tmp_path / "out")},
    }, ensure_ascii=False), encoding="utf-8")

    case_path = tmp_path / "case.json"
    case_path.write_text(json.dumps({
        "title": "身體五官",
        "grade": "國中七年級",
        "duration_minutes": 45,
        "natural_language_request": {
            "topic": "身體五官",
            "grade": "國中七年級",
            "outputs": ["exam", "slides", "interactive", "quiz"],
        },
        "vocabulary": [{"hanji": "身軀", "tailo_numeric": "sin1-khu1", "zh_tw": "身體"}],
        "dialogues": [{
            "role": "老師",
            "hanji": "咱來學身軀。",
            "tailo_numeric": "lan2 lai5 oh8 sin1-khu1.",
            "zh_tw": "我們來學身體。",
        }],
        "questions": [{
            "id": "q1",
            "question": "主題是啥？",
            "options": ["身軀", "天氣", "交通", "動物"],
            "answer_index": 0,
            "explanation": "本單元學身體。",
        }],
        "official_material_recommendations": {
            "interactive": [{"title": "官方素材", "resource_kind": "interactive", "page_url": "https://mhi.moe.edu.tw/"}]
        },
    }, ensure_ascii=False), encoding="utf-8")

    generator = MaterialGenerator(str(config_path))
    out_dir = tmp_path / "out"
    generator.generate_all(str(case_path), output_dir=str(out_dir), skip_media=True)

    manifest = {
        "request": {"topic": "身體五官", "grade": "國中七年級", "outputs": ["exam", "slides", "interactive", "quiz"]},
        "generation_options": {"skip_media": True},
        "official_material_recommendations": {"interactive": [{"title": "官方素材"}]},
        "outputs": {
            "lesson_structure": str(out_dir / "lesson_structure.json"),
            "student_worksheet": str(out_dir / "student_worksheet.docx"),
            "teacher_guide": str(out_dir / "teacher_guide.docx"),
            "exam_paper": str(out_dir / "exam_paper.docx"),
            "exam_answer_key": str(out_dir / "exam_answer_key.docx"),
            "quiz_bank": str(out_dir / "quiz_bank.json"),
            "quiz_teacher_key": str(out_dir / "quiz_teacher_key.md"),
            "slides": str(out_dir / "teaching_slides.pptx"),
            "interactive_website": str(out_dir / "interactive_website.html"),
            "teacher_review_report": str(out_dir / "teacher_review_report.md"),
            "video": None,
        },
    }
    (out_dir / "generation_manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    script_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "scripts", "validate_generation_output.py"))
    spec = importlib.util.spec_from_file_location("validate_generation_output", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    result = module.validate_output_folder(str(out_dir))
    assert result["ok"] is True
    module.write_reports(result, out_dir)
    assert (out_dir / "generation_validation.json").exists()
    assert (out_dir / "generation_validation.md").exists()


def test_generation_output_validator_rejects_failed_video(tmp_path):
    import importlib.util
    from docx import Document
    from pptx import Presentation

    out_dir = tmp_path
    lesson = {
        "title": "有機菜蔬",
        "vocabulary": [{"hanji": "菜蔬", "zh_tw": "蔬菜"}],
        "dialogues": [{"role": "老師", "hanji": "咱來學菜蔬。"}],
        "questions": [{
            "question": "主題是啥？",
            "options": ["菜蔬", "天氣"],
            "answer_index": 0,
            "explanation": "本單元學菜蔬。",
        }],
    }
    (out_dir / "lesson_structure.json").write_text(json.dumps(lesson, ensure_ascii=False), encoding="utf-8")
    for filename in ("student_worksheet.docx", "teacher_guide.docx", "exam_paper.docx", "exam_answer_key.docx"):
        doc = Document()
        doc.add_paragraph("有機菜蔬")
        doc.save(out_dir / filename)

    prs = Presentation()
    for _ in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(0, 0, 100, 100).text = "有機菜蔬"
    prs.save(out_dir / "teaching_slides.pptx")
    (out_dir / "interactive_website.html").write_text(
        "<h2>核心詞彙</h2><h2>情境會話</h2><h2>隨堂自我檢測</h2>",
        encoding="utf-8",
    )
    (out_dir / "teacher_review_report.md").write_text("有機菜蔬", encoding="utf-8")

    manifest = {
        "request": {"topic": "有機菜蔬", "grade": "國中七年級", "outputs": ["video"]},
        "generation_options": {"skip_media": False, "include_video": True},
        "video_generation": {
            "requested_in_text": True,
            "forced": False,
            "attempted": True,
            "success": False,
            "path": None,
            "skipped_reason": "generation_failed",
        },
        "official_material_recommendations": {"video": [{"title": "官方影片"}]},
        "outputs": {
            "lesson_structure": str(out_dir / "lesson_structure.json"),
            "student_worksheet": str(out_dir / "student_worksheet.docx"),
            "teacher_guide": str(out_dir / "teacher_guide.docx"),
            "exam_paper": str(out_dir / "exam_paper.docx"),
            "exam_answer_key": str(out_dir / "exam_answer_key.docx"),
            "slides": str(out_dir / "teaching_slides.pptx"),
            "interactive_website": str(out_dir / "interactive_website.html"),
            "teacher_review_report": str(out_dir / "teacher_review_report.md"),
            "video": None,
        },
    }
    (out_dir / "generation_manifest.json").write_text(json.dumps(manifest, ensure_ascii=False), encoding="utf-8")

    script_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "scripts", "validate_generation_output.py"))
    spec = importlib.util.spec_from_file_location("validate_generation_output", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    result = module.validate_output_folder(str(out_dir))
    assert result["ok"] is False
    assert any(issue["item"] == "video" and issue["severity"] == "error" for issue in result["issues"])


def test_project_goal_audit_uses_existing_evidence():
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "audit_project_goal.py"
    spec = importlib.util.spec_from_file_location("audit_project_goal", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    output_dir = project_root / "output" / "smoke_full_goal_all_outputs"
    report = module.audit_project_goal(output_dir)

    assert any(check["name"] == "official_repository_core_integrity" for check in report["checks"])
    assert any(check["name"] == "natural_language_requested_all_core_outputs" for check in report["checks"])
    assert any(check["name"] == "natural_language_video_success" for check in report["checks"])
    if output_dir.exists():
        assert report["ok"] is True


def test_formal_output_readiness_accepts_voxcpm_config(tmp_path, monkeypatch):
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "audit_formal_output_readiness.py"
    spec = importlib.util.spec_from_file_location("audit_formal_output_readiness", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    clone_script = tmp_path / "clone_batch.py"
    clone_script.write_text("# mock VoxCPM2 batch script\n", encoding="utf-8")
    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({
        "output": {"base_dir": str(tmp_path / "local_output")},
        "tts": {
            "provider": "voxcpm",
            "voxcpm": {
                "python": sys.executable,
                "script": str(clone_script),
                "voice": "三師爸台語",
            },
        },
        "piauim": {"provider": "ithuan"},
    }, ensure_ascii=False), encoding="utf-8")

    monkeypatch.setattr(module, "resolve_command", lambda name, prefer_cmd=False: f"C:/tools/{name}.exe")
    monkeypatch.setattr(module, "command_version", lambda command, args=None, timeout=8: {"ok": True, "output": "mock version"})

    report = module.audit_formal_output_readiness(config_path)

    assert report["ok"] is True
    assert any(check["name"] == "tts_provider_not_dummy" and check["ok"] for check in report["checks"])
    assert any(check["name"] == "voxcpm_script_exists" and check["ok"] for check in report["checks"])
    assert any(check["name"] == "tts_sample_generation" and check["severity"] == "warning" for check in report["checks"])


def test_formal_output_readiness_rejects_dummy_tts(tmp_path, monkeypatch):
    import importlib.util
    from pathlib import Path

    project_root = Path(__file__).resolve().parents[1]
    script_path = project_root / "scripts" / "audit_formal_output_readiness.py"
    spec = importlib.util.spec_from_file_location("audit_formal_output_readiness", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    config_path = tmp_path / "config.json"
    config_path.write_text(json.dumps({
        "output": {"base_dir": str(tmp_path / "local_output")},
        "tts": {"provider": "dummy"},
        "piauim": {"provider": "ithuan"},
    }, ensure_ascii=False), encoding="utf-8")

    monkeypatch.setattr(module, "resolve_command", lambda name, prefer_cmd=False: f"C:/tools/{name}.exe")
    monkeypatch.setattr(module, "command_version", lambda command, args=None, timeout=8: {"ok": True, "output": "mock version"})

    report = module.audit_formal_output_readiness(config_path)

    assert report["ok"] is False
    failed = [check["name"] for check in report["checks"] if check["severity"] == "error" and not check["ok"]]
    assert "tts_provider_not_dummy" in failed


def test_generation_output_validator_detects_topic_mismatch():
    import importlib.util

    script_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "scripts", "validate_generation_output.py"))
    spec = importlib.util.spec_from_file_location("validate_generation_output", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    issues = []
    result = module.validate_topic_alignment(
        {
            "vocabulary": [{"hanji": "食飯", "zh_tw": "吃飯"}],
            "dialogues": [{"hanji": "咱來去食飯。"}],
            "questions": [{"question": "食飯的意思是什麼？", "options": ["吃飯"]}],
        },
        {"request": {"topic": "身體五官"}},
        issues,
    )

    assert result["is_aligned"] is False
    assert any(issue["item"] == "topic_alignment" for issue in issues)


def test_vocabulary_db_tailo_consistency():
    """資料完整性：詞庫每筆 tailo_numeric 經轉換後，應與 tailo_diacritic 完全一致。
    （曾出現「偌濟錢」numeric 標 tsinn2 但 diacritic 為 tsînn 的不一致，會被反推覆蓋成錯誤聲調。）"""
    from rag.retriever import TaigiRetriever
    from tailo.validator import convert_sentence_numeric_to_diacritic
    retriever = TaigiRetriever()

    mismatches = []
    for entry in retriever.vocabulary_db:
        num = entry.get("tailo_numeric", "")
        dia = entry.get("tailo_diacritic", "")
        if not num or not dia:
            continue
        converted = convert_sentence_numeric_to_diacritic(num)
        if converted != dia:
            mismatches.append(f"{entry.get('hanji')}: {num} -> {converted} (期望 {dia})")

    assert not mismatches, "詞庫拼音不一致:\n" + "\n".join(mismatches)


def test_html_onclick_escapes_apostrophe(tmp_path):
    """漢字／音檔路徑含單引號時，不可破壞 onclick 的 JS 字串字面值。"""
    from generators.material_generator import MaterialGenerator
    gen = MaterialGenerator()
    data = {
        "title": "測試", "grade": "國中七年級", "duration_minutes": 45,
        "curriculum": {"learning_performance": [], "learning_content": []},
        "vocabulary": [{"hanji": "O'brien's 詞", "tailo_diacritic": "test",
                        "zh_tw": "測試", "audio_file": "", "image_file": ""}],
        "dialogues": [],
        "questions": []
    }
    gen._generate_html(data, str(tmp_path))
    html = (tmp_path / "interactive_website.html").read_text(encoding="utf-8")
    # 單引號必須被轉義成 \' 後嵌入，不可出現未轉義的破壞型片段
    assert "O\\'brien\\'s" in html


def test_vocabulary_db_expanded():
    from rag.retriever import TaigiRetriever
    retriever = TaigiRetriever()
    
    assert len(retriever.vocabulary_db) >= 30
    
    categories = {
        "食物": ["咖啡", "茶", "牛奶", "果子", "麵包"],
        "家人": ["阿公", "阿爸", "阿母", "小妹"],
        "學校": ["學校", "先生", "考試", "寫字"],
        "時間": ["今仔日", "明仔載", "昨昏"],
        "生活": ["你好", "再會", "電腦", "手機", "電影", "火車", "醫生", "公園"]
    }
    
    for cat, words in categories.items():
        for w in words:
            results = retriever.retrieve_vocabulary(w)
            assert len(results) > 0, f"找不到 {cat} 類詞彙: {w}"
            assert results[0]["hanji"] == w
            assert "pending" not in results[0].get("tailo_numeric", "")
